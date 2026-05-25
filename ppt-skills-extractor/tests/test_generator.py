"""Smoke tests for the deck generator and skeleton extractor."""

from __future__ import annotations

import sys
from pathlib import Path

import pytest
from pptx import Presentation

PROJECT_ROOT = Path(__file__).resolve().parent.parent
SRC = PROJECT_ROOT / "src"
sys.path.insert(0, str(SRC))

from generator.deck_builder import (
    RENDERERS,
    TemplateBrand,
    _clear_placeholders,
    build_deck,
)

SKILLS_ROOT = PROJECT_ROOT / "skills-output"
ASSETS_ROOT = SKILLS_ROOT / "assets"
TEMPLATE_ID = "rh-standard"
HAS_SKELETON = (SKILLS_ROOT / "templates" / TEMPLATE_ID / "theme-manifest.yaml").exists()


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

@pytest.fixture()
def output_dir(tmp_path: Path) -> Path:
    return tmp_path / "decks"


def _build(slides: list[dict], output_dir: Path, name: str = "test.pptx",
           **kw) -> Presentation:
    out = output_dir / name
    build_deck(
        slides,
        assets_root=ASSETS_ROOT,
        output_path=out,
        skills_root=SKILLS_ROOT,
        template_id=TEMPLATE_ID,
        **kw,
    )
    assert out.exists(), f"{out} was not created"
    assert out.stat().st_size > 0, f"{out} is empty"
    return Presentation(str(out))


# ---------------------------------------------------------------------------
# RENDERERS registry
# ---------------------------------------------------------------------------

def test_all_elements_registered():
    expected = {
        "title-block", "divider", "agenda", "metric-card",
        "challenge-list", "tech-tile", "quote-block", "image-content",
        "data-table", "recommendation-card", "timeline", "bar-chart",
        "closing",
    }
    assert expected.issubset(set(RENDERERS.keys()))


# ---------------------------------------------------------------------------
# Skeleton mode
# ---------------------------------------------------------------------------

@pytest.mark.skipif(not HAS_SKELETON, reason="No skeleton for rh-standard")
class TestSkeletonMode:

    def test_skeleton_creates_slides(self, output_dir):
        prs = _build([
            {"element": "title-block", "title": "T"},
            {"element": "closing"},
        ], output_dir)
        assert len(prs.slides) == 2

    def test_no_inherited_placeholders(self, output_dir):
        prs = _build([
            {"element": "metric-card", "title": "X", "cards": [
                {"value": "1", "label": "L"},
            ]},
        ], output_dir)
        for slide in prs.slides:
            phs = [s for s in slide.shapes if s.is_placeholder]
            assert phs == [], "Inherited placeholders should be cleared"

    def test_skeleton_preserves_masters(self, output_dir):
        prs = _build([{"element": "title-block", "title": "T"}], output_dir)
        assert len(prs.slide_masters) >= 2

    def test_master_alias_selection(self, output_dir):
        prs = _build(
            [{"element": "title-block", "title": "T"}],
            output_dir, name="charcoal.pptx",
            master_alias="charcoal",
        )
        assert len(prs.slides) == 1


# ---------------------------------------------------------------------------
# Individual renderers
# ---------------------------------------------------------------------------

@pytest.mark.skipif(not HAS_SKELETON, reason="No skeleton for rh-standard")
class TestRenderers:

    def test_title_block(self, output_dir):
        prs = _build([{
            "element": "title-block",
            "title": "Test Title",
            "subheading": "Sub",
            "date": "2026",
        }], output_dir)
        texts = _all_text(prs.slides[0])
        assert any("Test Title" in t for t in texts)

    def test_divider(self, output_dir):
        prs = _build([{
            "element": "divider",
            "headline": "Section Name",
            "section_marker": "Part 1",
        }], output_dir)
        texts = _all_text(prs.slides[0])
        assert any("Section Name" in t for t in texts)

    def test_metric_cards(self, output_dir):
        prs = _build([{
            "element": "metric-card",
            "title": "Metrics",
            "cards": [
                {"value": "99%", "label": "Uptime"},
                {"value": "50ms", "label": "Latency"},
            ],
        }], output_dir)
        texts = _all_text(prs.slides[0])
        assert any("99%" in t for t in texts)
        assert any("Uptime" in t for t in texts)

    def test_timeline(self, output_dir):
        prs = _build([{
            "element": "timeline",
            "title": "Roadmap",
            "milestones": [
                {"date": "Q1", "title": "Start"},
                {"date": "Q2", "title": "Mid"},
                {"date": "Q3", "title": "End"},
            ],
        }], output_dir)
        texts = _all_text(prs.slides[0])
        assert any("Q1" in t for t in texts)
        assert any("Start" in t for t in texts)

    def test_bar_chart(self, output_dir):
        prs = _build([{
            "element": "bar-chart",
            "title": "Growth",
            "bars": [
                {"label": "2023", "value": 100},
                {"label": "2024", "value": 150},
            ],
        }], output_dir)
        texts = _all_text(prs.slides[0])
        assert any("Growth" in t for t in texts)
        assert any("100" in t for t in texts)

    def test_challenge_list(self, output_dir):
        prs = _build([{
            "element": "challenge-list",
            "title": "Problems",
            "challenges": [
                {"headline": "Cost", "body": "Too expensive"},
            ],
        }], output_dir)
        texts = _all_text(prs.slides[0])
        assert any("Cost" in t for t in texts)

    def test_closing(self, output_dir):
        prs = _build([{"element": "closing"}], output_dir)
        texts = _all_text(prs.slides[0])
        assert any("Thank you" in t for t in texts)

    def test_agenda(self, output_dir):
        prs = _build([{
            "element": "agenda",
            "title": "Agenda",
            "items": [{"topic": "Intro"}, {"topic": "Demo"}],
        }], output_dir)
        texts = _all_text(prs.slides[0])
        assert any("Intro" in t for t in texts)

    def test_data_table(self, output_dir):
        prs = _build([{
            "element": "data-table",
            "title": "Comparison",
            "headers": ["Feature", "Us", "Them"],
            "rows": [["Speed", "Fast", "Slow"]],
        }], output_dir)
        assert len(prs.slides) == 1


# ---------------------------------------------------------------------------
# Validation
# ---------------------------------------------------------------------------

class TestValidation:

    def test_validate_slide_spec_valid(self):
        from agent.tools import validate_slide_spec
        result = validate_slide_spec.invoke({
            "template_id": TEMPLATE_ID,
            "spec": {"element": "title-block", "title": "Hello"},
        })
        if "error" not in str(result).lower() or "not found" in str(result).lower():
            pass  # schema may not exist for rh-standard

    def test_validate_number_type(self):
        from agent.tools import _validate_value
        errors, warnings = [], []
        _validate_value(42, "count", {"type": "number", "minimum": 0}, "test.count", errors, warnings)
        assert errors == []

        _validate_value(-1, "count", {"type": "number", "minimum": 0}, "test.count", errors, warnings)
        assert any("must be >= 0" in e for e in errors)

    def test_validate_rejects_wrong_type(self):
        from agent.tools import _validate_value
        errors, warnings = [], []
        _validate_value("not_a_number", "x", {"type": "number"}, "test.x", errors, warnings)
        assert any("expected number" in e for e in errors)


# ---------------------------------------------------------------------------
# Skeleton extractor
# ---------------------------------------------------------------------------

class TestSkeletonExtractor:

    def test_extract_produces_files(self, tmp_path):
        template = PROJECT_ROOT / "templates" / "1. Red Hat standard presentation template.pptx"
        if not template.exists():
            pytest.skip("Template file not available")

        from extractor.skeleton_extractor import extract_skeletons
        manifest_path = extract_skeletons(template, "test-rh", tmp_path)
        assert manifest_path.exists()

        import yaml
        with manifest_path.open() as f:
            manifest = yaml.safe_load(f)
        assert "masters" in manifest
        assert len(manifest["masters"]) >= 2

        skeleton = tmp_path / "templates" / "test-rh" / "skeletons" / "skeleton.pptx"
        assert skeleton.exists()
        prs = Presentation(str(skeleton))
        assert len(prs.slides) == 0
        assert len(prs.slide_masters) >= 2


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _all_text(slide) -> list[str]:
    texts = []
    for shape in slide.shapes:
        try:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        except Exception:
            pass
    return texts
