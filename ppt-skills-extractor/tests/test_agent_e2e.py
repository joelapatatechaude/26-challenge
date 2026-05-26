"""End-to-end tests for the Field Enablement PPT Agent."""

from __future__ import annotations

import asyncio
import json
import re
import sys
from pathlib import Path
from typing import Any

import pytest
from fastapi.testclient import TestClient
from langchain_core.messages import AIMessage, HumanMessage, SystemMessage
from pptx import Presentation

PROJECT_ROOT = Path(__file__).resolve().parent.parent
SRC = PROJECT_ROOT / "src"
sys.path.insert(0, str(SRC))

from agent.deck_types import (  # noqa: E402
    DECK_TYPES,
    default_outline,
    get_deck_type,
    list_deck_types,
)
from agent.graph import build_initial_state, create_building_graph, create_graph, create_planning_graph  # noqa: E402
from agent.prompts import COMPREHENSION_PROMPT  # noqa: E402
from agent.prompts import (  # noqa: E402
    CONTENT_WRITER_PROMPT,
    GEO_CONTEXT_PROMPT,
    OUTLINE_PLANNER_PROMPT,
    VALIDATOR_PROMPT,
)
from agent.tools import (  # noqa: E402
    build_presentation,
    get_element_schema,
    get_narrative_guide,
    list_elements,
    list_templates,
    search_icons,
    validate_slide_spec,
)

TEMPLATE_ID = "sales-enablement-2022"
EXPECTED_DECK_TYPES = {"elevator", "competitive", "power_hour", "questionnaire", "assessment"}


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------


@pytest.fixture()
def decks_dir(tmp_path, monkeypatch):
    """Redirect PPTX output to a temp directory and clean up after tests."""
    import agent.tools as tools_module

    monkeypatch.setattr(tools_module, "_PROJECT_ROOT", tmp_path)
    out = tmp_path / "decks"
    out.mkdir(parents=True, exist_ok=True)
    yield out
    for pptx in out.glob("*.pptx"):
        pptx.unlink(missing_ok=True)


def _minimal_slide_spec(element: str, topic: str = "Test Topic") -> dict[str, Any]:
    """Return a schema-valid minimal spec for common element types."""
    specs: dict[str, dict[str, Any]] = {
        "title-block": {
            "element": "title-block",
            "title": topic,
            "subheading": "Field Enablement",
        },
        "image-content": {
            "element": "image-content",
            "title": "Overview",
            "bullets": [{"body": "Key sovereignty benefit for customers."}],
        },
        "challenge-list": {
            "element": "challenge-list",
            "title": "Customer Challenges",
            "challenges": [
                {"headline": "Compliance", "body": "Meeting regional data laws."},
                {"headline": "Control", "body": "Retaining operational sovereignty."},
                {"headline": "Cost", "body": "Avoiding vendor lock-in."},
            ],
        },
        "recommendation-card": {
            "element": "recommendation-card",
            "title": "Next Steps",
            "cards": [
                {"headline": "Assess", "body": "Run a sovereignty readiness workshop."},
                {"headline": "Pilot", "body": "Deploy a hybrid cloud proof of concept."},
            ],
        },
        "closing": {"element": "closing"},
        "agenda": {
            "element": "agenda",
            "title": "Agenda",
            "items": [{"topic": "Intro"}, {"topic": "Solutions"}],
        },
        "divider": {
            "element": "divider",
            "headline": "Section",
            "section_marker": "01",
        },
        "metric-card": {
            "element": "metric-card",
            "title": "Metrics",
            "cards": [{"value": "99%", "label": "Uptime"}],
        },
        "tech-tile": {
            "element": "tech-tile",
            "title": "Platform",
            "tiles": [{"headline": "Open Hybrid Cloud", "body": "Flexible deployment."}],
        },
        "data-table": {
            "element": "data-table",
            "title": "Comparison",
            "headers": ["Feature", "Red Hat"],
            "rows": [["Open source", "Yes"]],
        },
    }
    if element in specs:
        return specs[element]
    return {"element": element, "title": topic}


class MockLLM:
    """Deterministic LLM stub — never calls external APIs."""

    def __init__(self, deck_type: str = "elevator", topic: str = "Digital Sovereignty"):
        self.deck_type = deck_type
        self.topic = topic
        self.calls: list[list] = []

    async def ainvoke(self, messages, config=None):
        self.calls.append(messages)
        sys_text = " ".join(
            str(m.content)
            for m in messages
            if isinstance(m, (SystemMessage,)) or getattr(m, "type", "") == "system"
        )
        user_text = " ".join(
            str(m.content) for m in messages if isinstance(m, HumanMessage)
        )
        combined = f"{sys_text}\n{user_text}"

        if GEO_CONTEXT_PROMPT[:40] in sys_text or "sovereignty expert" in sys_text.lower():
            return AIMessage(content="- GDPR and Schrems II drive EU data residency\n- Hybrid cloud supports sovereignty")

        if OUTLINE_PLANNER_PROMPT[:40] in user_text or "building a presentation outline" in user_text.lower():
            outline = default_outline(self.deck_type)
            return AIMessage(content=json.dumps(outline))

        if CONTENT_WRITER_PROMPT[:40] in sys_text or "slide content writer" in sys_text.lower():
            match = re.search(r'"element"\s*:\s*"([^"]+)"', user_text)
            element = match.group(1) if match else "title-block"
            spec = _minimal_slide_spec(element, self.topic)
            return AIMessage(content=json.dumps(spec))

        if VALIDATOR_PROMPT[:30] in combined or "QA validator" in combined:
            return AIMessage(content="VALIDATED")

        if COMPREHENSION_PROMPT[:40] in sys_text or "summarising what a user wants" in sys_text.lower():
            return AIMessage(
                content=json.dumps(
                    {
                        "deck_mode": "fresh",
                        "customer": "",
                        "audience": "Field sellers",
                        "geo_context": "",
                        "document_ref": "no documents",
                        "summary": "Planning deck from prompt. Key themes will be drawn from your prompt.",
                        "themes_status": "to_be_extracted",
                        "gaps": [],
                    }
                )
            )

        return AIMessage(content="OK")


@pytest.fixture()
def mock_llm():
    return MockLLM()


@pytest.fixture()
def mock_graph(mock_llm):
    return create_graph(mock_llm)


@pytest.fixture()
def mock_planning_graph(mock_llm):
    return create_planning_graph(mock_llm)


@pytest.fixture()
def api_client(mock_llm):
    import agent.api as api_module

    api_module._graph = create_graph(mock_llm)
    api_module._planning_graph = create_planning_graph(mock_llm)
    api_module._building_graph = create_building_graph(mock_llm)
    with TestClient(api_module.app) as client:
        yield client
    api_module._graph = None
    api_module._planning_graph = None
    api_module._building_graph = None


# ---------------------------------------------------------------------------
# A1. Deck Types Registry
# ---------------------------------------------------------------------------


class TestDeckTypesRegistry:
    def test_all_five_deck_types_registered(self):
        assert set(DECK_TYPES.keys()) == EXPECTED_DECK_TYPES

    @pytest.mark.parametrize("deck_type", sorted(EXPECTED_DECK_TYPES))
    def test_each_deck_type_has_valid_config(self, deck_type):
        spec = DECK_TYPES[deck_type]
        assert spec.id == deck_type
        assert spec.slide_count_min > 0
        assert spec.slide_count_max >= spec.slide_count_min
        assert len(spec.elements) >= 2
        assert len(spec.elements) <= spec.slide_count_max + 2
        assert spec.elements[0] == "title-block"
        assert spec.elements[-1] == "closing"
        assert spec.pptx_filename
        assert spec.pptx_filename.lower().endswith((".pptx", ".pptx".upper()))

    @pytest.mark.parametrize(
        "deck_type,expected_name",
        [
            ("elevator", "5-Minute Elevator Pitch"),
            ("competitive", "15-Minute Competitive Conversation"),
            ("power_hour", "Power Hour Deep Dive"),
            ("questionnaire", "Customer Discovery Questionnaire"),
            ("assessment", "Readiness Assessment"),
        ],
    )
    def test_get_deck_type_returns_correct_config(self, deck_type, expected_name):
        spec = get_deck_type(deck_type)
        assert spec.id == deck_type
        assert spec.name == expected_name

    def test_get_deck_type_invalid_falls_back_to_competitive(self):
        spec = get_deck_type("not-a-real-type")
        assert spec.id == "competitive"

    def test_get_deck_type_normalizes_hyphens(self):
        assert get_deck_type("power-hour").id == "power_hour"

    def test_list_deck_types_returns_all_five(self):
        types = list_deck_types()
        ids = {t["id"] for t in types}
        assert ids == EXPECTED_DECK_TYPES
        for entry in types:
            assert "slide_count_min" in entry
            assert "slide_count_max" in entry
            assert "element_count" in entry
            assert "pptx_filename" in entry


# ---------------------------------------------------------------------------
# A2. Tools
# ---------------------------------------------------------------------------


class TestTools:
    def test_list_templates_returns_nonempty_with_expected_keys(self):
        templates = list_templates.invoke({})
        assert templates
        first = templates[0]
        assert "template_id" in first
        assert "name" in first
        assert "primary_color" in first
        assert "heading_font" in first
        assert "body_font" in first
        assert any(t["template_id"] == TEMPLATE_ID for t in templates)

    def test_list_elements_returns_id_layout_purpose(self):
        elements = list_elements.invoke({"template_id": TEMPLATE_ID})
        assert elements
        title = next(e for e in elements if e["id"] == "title-block")
        assert title["layout"]
        assert title["purpose"]

    def test_get_element_schema_returns_fields(self):
        schema = get_element_schema.invoke(
            {"template_id": TEMPLATE_ID, "element": "title-block"}
        )
        assert schema["element"] == "title-block"
        assert "fields" in schema
        assert "title" in schema["fields"]
        assert schema["fields"]["element"]["const"] == "title-block"

    def test_search_icons_security_returns_tags(self):
        results = search_icons.invoke({"query": "security"})
        assert results
        assert all("tags" in r for r in results)
        assert any(
            r.get("category") == "security"
            or any("security" in str(tag).lower() for tag in r.get("tags", []))
            or "security" in str(r.get("id", "")).lower()
            for r in results
        )

    def test_get_narrative_guide_returns_nonempty_string(self):
        guide = get_narrative_guide.invoke({"template_id": TEMPLATE_ID})
        assert isinstance(guide, str)
        assert len(guide.strip()) > 20

    def test_validate_slide_spec_valid(self):
        result = validate_slide_spec.invoke(
            {
                "template_id": TEMPLATE_ID,
                "spec": {"element": "title-block", "title": "Valid Title"},
            }
        )
        assert result["valid"] is True
        assert result["errors"] == []

    def test_validate_slide_spec_invalid(self):
        result = validate_slide_spec.invoke(
            {
                "template_id": TEMPLATE_ID,
                "spec": {"element": "title-block"},
            }
        )
        assert result["valid"] is False
        assert any("title" in err for err in result["errors"])

    def test_build_presentation_minimal_two_slides(self, decks_dir):
        slides = [
            {"element": "title-block", "title": "E2E Test Deck"},
            {"element": "closing"},
        ]
        result = build_presentation.invoke(
            {
                "template_id": TEMPLATE_ID,
                "slides": slides,
                "output_name": "e2e-minimal.pptx",
            }
        )
        assert "Validation failed" not in result
        assert "Error" not in result.split(":")[0]
        out = decks_dir / "e2e-minimal.pptx"
        assert out.exists()
        assert out.stat().st_size > 0
        prs = Presentation(str(out))
        assert len(prs.slides) == 2


# ---------------------------------------------------------------------------
# A3. API Endpoints
# ---------------------------------------------------------------------------


class TestAPIEndpoints:
    def test_health_returns_ok(self, api_client):
        resp = api_client.get("/health")
        assert resp.status_code == 200
        body = resp.json()
        assert body["status"] == "ok"

    def test_deck_types_returns_all_five(self, api_client):
        resp = api_client.get("/api/v1/deck-types")
        assert resp.status_code == 200
        ids = {d["id"] for d in resp.json()}
        assert ids == EXPECTED_DECK_TYPES

    def test_generate_missing_topic_returns_422(self, api_client):
        resp = api_client.post("/api/v1/generate", json={"deck_type": "elevator"})
        assert resp.status_code == 422

    def test_generate_empty_topic_returns_422(self, api_client):
        resp = api_client.post(
            "/api/v1/generate",
            json={"topic": "   ", "deck_type": "elevator"},
        )
        assert resp.status_code == 422


# ---------------------------------------------------------------------------
# B4. Graph Flow with Mocked LLM
# ---------------------------------------------------------------------------


class TestPlanningGraphComprehension:
    @pytest.mark.asyncio
    async def test_planning_graph_emits_comprehension_event(
        self, mock_planning_graph, monkeypatch
    ):
        monkeypatch.setenv("PLANNING_COMPREHENSION", "true")
        import agent.graph as graph_module

        monkeypatch.setattr(graph_module, "_PLANNING_COMPREHENSION", True)

        initial = build_initial_state(
            topic="Planning comprehension test",
            deck_type="elevator",
            template_id=TEMPLATE_ID,
        )
        event_queue: asyncio.Queue = asyncio.Queue()
        config = {"configurable": {"event_queue": event_queue}}

        final = await mock_planning_graph.ainvoke(initial, config)

        events: list[dict] = []
        while not event_queue.empty():
            item = event_queue.get_nowait()
            if item is not None:
                events.append(item)

        event_names = [e["event"] for e in events]
        assert "comprehension" in event_names
        assert final.get("outline")
        assert final.get("intent_summary") is not None
        comprehension_idx = event_names.index("comprehension")
        progress_after = [
            i for i, name in enumerate(event_names) if name == "progress"
        ]
        outline_progress = [
            i
            for i, name in enumerate(event_names)
            if name == "progress"
            and "outline" in json.loads(events[i]["data"]).get("step", "").lower()
        ]
        if outline_progress:
            assert comprehension_idx < outline_progress[-1]

    @pytest.mark.asyncio
    async def test_uploaded_files_description_forwarded_to_deck_state(self):
        state = build_initial_state(
            topic="Upload desc test",
            uploaded_files_description="brand.pptx (theme template), report.pdf (reference)",
        )
        assert state["uploaded_files_description"] == (
            "brand.pptx (theme template), report.pdf (reference)"
        )
        assert state.get("intent_summary") is None

    def test_generate_request_forwards_uploaded_files_description(self, api_client):
        resp = api_client.post(
            "/api/v1/generate",
            json={
                "topic": "SSE planning test",
                "deck_type": "elevator",
                "template_id": TEMPLATE_ID,
                "uploaded_files_description": "brand.pptx (theme template)",
            },
        )
        assert resp.status_code == 200
        events = _parse_sse_events(resp.text)
        event_names = [e[0] for e in events]
        assert "comprehension" in event_names
        outline_idx = event_names.index("outline_ready")
        comprehension_idx = event_names.index("comprehension")
        assert comprehension_idx < outline_idx


class TestGraphFlowMockedLLM:
    @pytest.mark.asyncio
    async def test_full_graph_flow_produces_pptx(self, mock_graph, decks_dir):
        topic = "Sovereignty E2E Test"
        initial = build_initial_state(
            topic=topic,
            deck_type="elevator",
            template_id=TEMPLATE_ID,
        )
        event_queue: asyncio.Queue = asyncio.Queue()
        config = {"configurable": {"event_queue": event_queue}}

        final = await mock_graph.ainvoke(initial, config)

        assert final.get("outline")
        assert len(final["outline"]) == 5
        assert len(final["slide_specs"]) == 5
        assert final["deck_path"]
        deck_path = Path(final["deck_path"])
        assert deck_path.exists()
        assert deck_path.suffix.lower() == ".pptx"
        prs = Presentation(str(deck_path))
        assert len(prs.slides) == 5

        queued_events: list[dict] = []
        while not event_queue.empty():
            item = event_queue.get_nowait()
            if item is not None:
                queued_events.append(item)

        event_types = {e["event"] for e in queued_events}
        assert "progress" in event_types
        assert "slide_spec" in event_types
        assert "deck_ready" in event_types

        progress = final.get("progress_events") or []
        assert progress
        progress_types = {p["event"] for p in progress}
        assert "slide_spec" in progress_types or "deck_ready" in progress_types


# ---------------------------------------------------------------------------
# B5. SSE Streaming with Mocked LLM
# ---------------------------------------------------------------------------


def _parse_sse_events(raw: str) -> list[tuple[str, dict]]:
    """Parse SSE text into (event_type, data_dict) pairs."""
    events: list[tuple[str, dict]] = []
    current_event = "message"
    for line in raw.splitlines():
        if line.startswith("event:"):
            current_event = line.split(":", 1)[1].strip()
        elif line.startswith("data:"):
            payload = line.split(":", 1)[1].strip()
            events.append((current_event, json.loads(payload)))
    return events


class TestSSEStreaming:
    def test_generate_streams_planning_events_in_order(self, api_client):
        resp = api_client.post(
            "/api/v1/generate",
            json={
                "topic": "SSE Stream Test",
                "deck_type": "elevator",
                "template_id": TEMPLATE_ID,
            },
        )
        assert resp.status_code == 200
        assert "text/event-stream" in resp.headers.get("content-type", "")

        events = _parse_sse_events(resp.text)
        assert events

        event_types = [e[0] for e in events]
        assert event_types[0] == "progress"
        assert "comprehension" in event_types
        assert "outline_ready" in event_types
        assert event_types[-1] == "completed"

        comprehension_idx = event_types.index("comprehension")
        outline_ready_idx = event_types.index("outline_ready")
        completed_idx = event_types.index("completed")
        assert comprehension_idx < outline_ready_idx < completed_idx

        for event_type, data in events:
            assert isinstance(data, dict)
            if event_type == "progress":
                assert "status" in data
            elif event_type == "comprehension":
                assert data.get("themes_status") == "to_be_extracted"
            elif event_type == "outline_ready":
                assert "outline" in data
                assert "job_id" in data
            elif event_type == "completed":
                assert data.get("status") in ("outline_ready", "completed", "failed")


# ---------------------------------------------------------------------------
# C6. Tool Integration — build deck without LLM
# ---------------------------------------------------------------------------


class TestBuildDeckWithoutLLM:
    def test_hardcoded_elevator_spec(self, decks_dir):
        """Mimic content_writer output for a full elevator deck."""
        slides = [
            _minimal_slide_spec(el["element"], "Direct Build Test")
            for el in default_outline("elevator")
        ]
        output_name = "direct-elevator.pptx"
        result = build_presentation.invoke(
            {
                "template_id": TEMPLATE_ID,
                "slides": slides,
                "output_name": output_name,
            }
        )
        assert "Validation failed" not in result

        out = decks_dir / output_name
        assert out.exists()
        prs = Presentation(str(out))
        assert len(prs.slides) == len(slides)
