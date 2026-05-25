"""Parse a PPTX file into RawSlide / RawShape Pydantic objects."""

from __future__ import annotations

import hashlib
import subprocess
import tempfile
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu
from rich.console import Console

from schemas import RawShape, RawSlide

console = Console()

EMU_PER_INCH = 914400


def _hex(color) -> Optional[str]:
    """Convert a pptx RGBColor to a #RRGGBB hex string, or None."""
    try:
        if color is None:
            return None
        return f"#{color.rgb}"
    except Exception:
        return None


def _shape_fill_color(shape) -> Optional[str]:
    try:
        fill = shape.fill
        if fill.type is None:
            return None
        from pptx.enum.dml import MSO_THEME_COLOR
        if fill.fore_color and fill.fore_color.rgb:
            return f"#{fill.fore_color.rgb}"
    except Exception:
        pass
    return None


def _shape_line_color(shape) -> Optional[str]:
    try:
        line = shape.line
        if line and line.color and line.color.rgb:
            return f"#{line.color.rgb}"
    except Exception:
        pass
    return None


def _shape_line_width(shape) -> Optional[float]:
    try:
        line = shape.line
        if line and line.width:
            return line.width.pt
    except Exception:
        pass
    return None


def _collect_text_runs(shape) -> tuple[str, list[str], list[str]]:
    """Return (full_text, bold_runs, italic_runs) from a shape's text frame."""
    full_text = ""
    bold_runs: list[str] = []
    italic_runs: list[str] = []
    try:
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text_frame:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                t = run.text
                                if not t:
                                    continue
                                full_text += t
                                if run.font.bold:
                                    bold_runs.append(t)
                                if run.font.italic:
                                    italic_runs.append(t)
                            full_text += "\n"
        else:
            tf = shape.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    t = run.text
                    if not t:
                        continue
                    full_text += t
                    if run.font.bold:
                        bold_runs.append(t)
                    if run.font.italic:
                        italic_runs.append(t)
                full_text += "\n"
    except Exception:
        try:
            full_text = shape.text
        except Exception:
            pass
    return full_text.strip(), bold_runs, italic_runs


def _font_sizes(shape) -> list[float]:
    sizes: list[float] = []
    try:
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text_frame:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font.size:
                                    sizes.append(run.font.size.pt)
        else:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        sizes.append(run.font.size.pt)
    except Exception:
        pass
    return sizes


def _font_names_and_colors(shape) -> tuple[list[str], list[str]]:
    names: list[str] = []
    colors: list[str] = []
    try:
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text_frame:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font.name and run.font.name not in names:
                                    names.append(run.font.name)
                                try:
                                    if run.font.color and run.font.color.rgb:
                                        hex_c = f"#{run.font.color.rgb}"
                                        if hex_c not in colors:
                                            colors.append(hex_c)
                                except Exception:
                                    pass
        else:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name and run.font.name not in names:
                        names.append(run.font.name)
                    try:
                        if run.font.color and run.font.color.rgb:
                            hex_c = f"#{run.font.color.rgb}"
                            if hex_c not in colors:
                                colors.append(hex_c)
                    except Exception:
                        pass
    except Exception:
        pass
    return names, colors


def _image_hash(shape) -> Optional[str]:
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            blob = shape.image.blob
            return hashlib.sha1(blob).hexdigest()
    except Exception:
        pass
    return None


def parse_pptx(path: Path) -> list[RawSlide]:
    """
    Parse all slides in a PPTX file into RawSlide objects.
    If the file is ODP, it is first converted to PPTX via LibreOffice headless.
    """
    path = _ensure_pptx(path)
    prs = Presentation(str(path))
    slide_w = prs.slide_width or Emu(9144000)
    slide_h = prs.slide_height or Emu(6858000)

    raw_slides: list[RawSlide] = []

    for idx, slide in enumerate(prs.slides):
        layout_name = ""
        try:
            layout_name = slide.slide_layout.name
        except Exception:
            pass

        title_text = ""
        try:
            title_text = slide.shapes.title.text if slide.shapes.title else ""
        except Exception:
            pass

        shapes: list[RawShape] = []
        for shape in slide.shapes:
            left_pct = round(shape.left / slide_w * 100, 2) if shape.left else 0.0
            top_pct = round(shape.top / slide_h * 100, 2) if shape.top else 0.0
            width_pct = round(shape.width / slide_w * 100, 2) if shape.width else 0.0
            height_pct = round(shape.height / slide_h * 100, 2) if shape.height else 0.0

            text, bold_runs, italic_runs = _collect_text_runs(shape)
            font_sizes = _font_sizes(shape)
            font_names, font_colors = _font_names_and_colors(shape)

            raw_shapes = RawShape(
                shape_id=shape.shape_id,
                shape_type=str(shape.shape_type),
                name=shape.name,
                left_pct=left_pct,
                top_pct=top_pct,
                width_pct=width_pct,
                height_pct=height_pct,
                text=text,
                font_sizes=font_sizes,
                font_names=font_names,
                font_colors=font_colors,
                fill_color=_shape_fill_color(shape),
                line_color=_shape_line_color(shape),
                line_width_pt=_shape_line_width(shape),
                bold_runs=bold_runs,
                italic_runs=italic_runs,
                image_hash=_image_hash(shape),
            )
            shapes.append(raw_shapes)

        raw_slides.append(RawSlide(
            file=path.name,
            slide_index=idx,
            layout_name=layout_name,
            title=title_text,
            shapes=shapes,
        ))

    console.print(f"  [green]Parsed[/green] {len(raw_slides)} slides from [bold]{path.name}[/bold]")
    return raw_slides


def _ensure_pptx(path: Path) -> Path:
    """If the file is ODP, convert to PPTX via LibreOffice headless."""
    if path.suffix.lower() != ".odp":
        return path
    console.print(f"  [yellow]ODP detected — converting via LibreOffice...[/yellow]")
    with tempfile.TemporaryDirectory() as tmpdir:
        result = subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "pptx", str(path), "--outdir", tmpdir],
            capture_output=True, text=True
        )
        if result.returncode != 0:
            raise RuntimeError(f"LibreOffice conversion failed:\n{result.stderr}")
        converted = Path(tmpdir) / (path.stem + ".pptx")
        if not converted.exists():
            raise FileNotFoundError(f"Expected converted file at {converted}")
        dest = path.with_suffix(".pptx")
        dest.write_bytes(converted.read_bytes())
        console.print(f"  [green]Converted to[/green] {dest.name}")
        return dest
