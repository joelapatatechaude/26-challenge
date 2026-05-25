"""LangChain tools for discovering templates, querying schemas, and building decks."""

from __future__ import annotations

import os
import sys
from pathlib import Path
from typing import Any

import yaml
from langchain_core.tools import tool

_PROJECT_ROOT = Path(__file__).resolve().parent.parent.parent
_SRC = _PROJECT_ROOT / "src"


def _default_skills_root() -> Path:
    env = os.environ.get("SKILLS_ROOT")
    if env:
        p = Path(env)
        return p if p.is_absolute() else _PROJECT_ROOT / p
    return _PROJECT_ROOT / "skills-output"


SKILLS_ROOT: Path = _default_skills_root()


def _get_build_deck():
    """Import deck_builder using PYTHONPATH=src layout."""
    if str(_SRC) not in sys.path:
        sys.path.insert(0, str(_SRC))
    from generator.deck_builder import build_deck

    return build_deck


def _load_yaml(path: Path) -> dict[str, Any]:
    if not path.exists():
        return {}
    with path.open() as f:
        return yaml.safe_load(f) or {}


def _template_dir(template_id: str) -> Path:
    return SKILLS_ROOT / "templates" / template_id


def _agent_index_path(template_id: str) -> Path:
    return _template_dir(template_id) / "agent_index.yaml"


def _load_agent_index(template_id: str) -> dict[str, Any]:
    return _load_yaml(_agent_index_path(template_id))


def _brand_from_template(template_id: str) -> dict[str, Any]:
    data = _load_yaml(_template_dir(template_id) / "template.yaml")
    theme = data.get("theme", {})
    colors = theme.get("colors", {})
    return {
        "primary_color": theme.get("primary_color", "#EE0000"),
        "text_color": colors.get("dk1", "#151515"),
        "heading_font": theme.get("font_heading", "Red Hat Display"),
        "body_font": theme.get("font_body", "Red Hat Text"),
    }


def _schema_has_icon_ref(schema: Any) -> bool:
    if isinstance(schema, dict):
        if schema.get("type") == "icon_ref":
            return True
        return any(_schema_has_icon_ref(v) for v in schema.values())
    if isinstance(schema, list):
        return any(_schema_has_icon_ref(item) for item in schema)
    return False


def _fields_from_schema(fields: dict[str, Any]) -> tuple[list[str], list[str], bool]:
    required: list[str] = []
    optional: list[str] = []
    for name, spec in fields.items():
        if name == "element" or not isinstance(spec, dict):
            continue
        if spec.get("required"):
            required.append(name)
        else:
            optional.append(name)
    return required, optional, _schema_has_icon_ref(fields)


def _element_entry_from_yaml(elem_id: str, data: dict[str, Any]) -> dict[str, Any]:
    required, optional, supports_icons = _fields_from_schema(data.get("fields", {}))
    return {
        "id": elem_id,
        "layout": data.get("blueprint_layout", ""),
        "purpose": data.get("description", ""),
        "required_fields": required,
        "optional_fields": optional,
        "supports_icons": supports_icons,
    }


def _list_elements_fallback(template_id: str) -> list[dict[str, Any]]:
    elements_dir = _template_dir(template_id) / "elements"
    if not elements_dir.exists():
        return []
    entries: list[dict[str, Any]] = []
    for path in sorted(elements_dir.glob("*.yaml")):
        data = _load_yaml(path)
        if data:
            entries.append(_element_entry_from_yaml(path.stem, data))
    return entries


def _icon_search_score(icon: dict[str, Any], query: str) -> int:
    q = query.strip().lower()
    if not q:
        return 0

    category = (icon.get("category") or "").lower()
    tags = [str(t).lower() for t in icon.get("tags", [])]
    name = (icon.get("name") or "").lower()
    description = (icon.get("description") or "").lower()
    icon_id = (icon.get("id") or "").lower()

    best = 0
    tokens = q.split()
    for token in tokens:
        if category == token:
            best = max(best, 100)
        elif token in category:
            best = max(best, 85)
        if any(tag == token for tag in tags):
            best = max(best, 75)
        elif any(token in tag for tag in tags):
            best = max(best, 60)
        if token in name:
            best = max(best, 70)
        if token in icon_id:
            best = max(best, 65)
        if token in description:
            best = max(best, 35)
    return best


def _validate_value(
    value: Any,
    field_name: str,
    spec: dict[str, Any],
    path: str,
    errors: list[str],
    warnings: list[str],
) -> None:
    field_type = spec.get("type", "string")

    if field_type == "string":
        if not isinstance(value, str):
            errors.append(f"{path}: expected string, got {type(value).__name__}")
            return
        const = spec.get("const")
        if const is not None and value != const:
            errors.append(f"{path}: must be {const!r}, got {value!r}")
        max_len = spec.get("max_length")
        if max_len is not None and len(value) > max_len:
            errors.append(f"{path}: exceeds max_length {max_len} ({len(value)} chars)")
        return

    if field_type == "icon_ref":
        if not isinstance(value, str):
            errors.append(f"{path}: expected icon asset_id string, got {type(value).__name__}")
            return
        catalog = _load_yaml(SKILLS_ROOT / "assets" / "icon_catalog.yaml")
        known_ids = {icon.get("id") for icon in catalog.get("icons", [])}
        if known_ids and value not in known_ids:
            warnings.append(f"{path}: icon {value!r} not found in icon_catalog.yaml")
        return

    if field_type == "image_ref":
        if not isinstance(value, str):
            errors.append(f"{path}: expected image asset_id string, got {type(value).__name__}")
            return
        catalog = _load_yaml(SKILLS_ROOT / "assets" / "image_catalog.yaml")
        known_ids = {img.get("id") for img in catalog.get("images", [])}
        if known_ids and value not in known_ids:
            warnings.append(f"{path}: image {value!r} not found in image_catalog.yaml")
        return

    if field_type == "array":
        if not isinstance(value, list):
            errors.append(f"{path}: expected array, got {type(value).__name__}")
            return
        min_items = spec.get("min_items")
        max_items = spec.get("max_items")
        if min_items is not None and len(value) < min_items:
            errors.append(f"{path}: requires at least {min_items} items, got {len(value)}")
        if max_items is not None and len(value) > max_items:
            errors.append(f"{path}: allows at most {max_items} items, got {len(value)}")
        item_schema = spec.get("item_schema")
        if item_schema:
            for i, item in enumerate(value):
                if not isinstance(item, dict):
                    errors.append(f"{path}[{i}]: expected object, got {type(item).__name__}")
                    continue
                _validate_object(item, item_schema, f"{path}[{i}]", errors, warnings)
        return

    if field_type == "number":
        if not isinstance(value, (int, float)):
            errors.append(f"{path}: expected number, got {type(value).__name__}")
            return
        minimum = spec.get("minimum")
        if minimum is not None and value < minimum:
            errors.append(f"{path}: must be >= {minimum}, got {value}")
        maximum = spec.get("maximum")
        if maximum is not None and value > maximum:
            errors.append(f"{path}: must be <= {maximum}, got {value}")
        return

    if not isinstance(value, str):
        errors.append(f"{path}: unsupported type {field_type!r} for value {type(value).__name__}")


def _validate_object(
    data: dict[str, Any],
    schema: dict[str, Any],
    path: str,
    errors: list[str],
    warnings: list[str],
) -> None:
    for name, spec in schema.items():
        if not isinstance(spec, dict):
            continue
        present = name in data
        if spec.get("required") and not present:
            errors.append(f"{path}.{name}: required field missing")
            continue
        if not present:
            continue
        _validate_value(data[name], name, spec, f"{path}.{name}", errors, warnings)


def _validate_spec_against_schema(spec: dict[str, Any], fields: dict[str, Any]) -> dict[str, Any]:
    errors: list[str] = []
    warnings: list[str] = []
    _validate_object(spec, fields, "spec", errors, warnings)
    return {"valid": not errors, "errors": errors, "warnings": warnings}


@tool
def list_templates() -> list[dict]:
    """List available presentation templates with brand summaries.

    Returns template ID, name, primary color, fonts for each template."""
    index = _load_yaml(SKILLS_ROOT / "index.yaml")
    results: list[dict] = []
    for tmpl in index.get("templates", []):
        template_id = tmpl.get("template_id", "")
        agent = _load_agent_index(template_id)
        brand = agent.get("brand") or _brand_from_template(template_id)
        results.append({
            "template_id": template_id,
            "name": agent.get("name") or tmpl.get("name", template_id),
            "primary_color": brand.get("primary_color"),
            "heading_font": brand.get("heading_font"),
            "body_font": brand.get("body_font"),
        })
    return results


@tool
def get_element_schema(template_id: str, element: str) -> dict:
    """Get the full field schema for a slide element type.

    Use this to understand what fields are needed to create a specific slide type."""
    path = _template_dir(template_id) / "elements" / f"{element}.yaml"
    if not path.exists():
        return {"error": f"Element {element!r} not found for template {template_id!r}"}
    data = _load_yaml(path)
    return {
        "element": element,
        "template_id": template_id,
        "fields": data.get("fields", {}),
        "constraints": data.get("constraints", []),
        "when_to_use": data.get("when_to_use", []),
    }


@tool
def list_elements(template_id: str) -> list[dict]:
    """List all available slide element types for a template.

    Returns element ID, purpose, required/optional fields, layout, and icon support."""
    agent = _load_agent_index(template_id)
    if agent.get("elements"):
        return agent["elements"]
    return _list_elements_fallback(template_id)


@tool
def search_icons(query: str, top_k: int = 5) -> list[dict]:
    """Search icons by semantic tags or category.

    Query can be a category name (security, cloud) or descriptive tags (shield, network)."""
    catalog = _load_yaml(SKILLS_ROOT / "assets" / "icon_catalog.yaml")
    scored: list[tuple[int, dict]] = []
    for icon in catalog.get("icons", []):
        score = _icon_search_score(icon, query)
        if score > 0:
            scored.append((score, icon))

    scored.sort(key=lambda pair: (-pair[0], pair[1].get("id", "")))
    return [
        {
            "id": icon.get("id"),
            "file": icon.get("file"),
            "category": icon.get("category"),
            "tags": icon.get("tags", []),
            "description": icon.get("description"),
            "suggested_elements": icon.get("suggested_elements", []),
            "relevance": score,
        }
        for score, icon in scored[:top_k]
    ]


@tool
def search_images(query: str, top_k: int = 5) -> list[dict]:
    """Search stock images by semantic tags or category.

    Query can be a category name (business, urban) or descriptive tags (city, skyline)."""
    catalog = _load_yaml(SKILLS_ROOT / "assets" / "image_catalog.yaml")
    scored: list[tuple[int, dict]] = []
    for image in catalog.get("images", []):
        score = _icon_search_score(image, query)
        if score > 0:
            scored.append((score, image))

    scored.sort(key=lambda pair: (-pair[0], pair[1].get("id", "")))
    return [
        {
            "id": image.get("id"),
            "file": image.get("file"),
            "category": image.get("category"),
            "tags": image.get("tags", []),
            "description": image.get("description"),
            "relevance": score,
        }
        for score, image in scored[:top_k]
    ]


@tool
def get_narrative_guide(template_id: str) -> str:
    """Get the recommended deck narrative structure for a template.

    Returns guidance on slide ordering, section structure, and best practices."""
    agent = _load_agent_index(template_id)
    guide = agent.get("narrative_guide")
    if guide:
        return guide if isinstance(guide, str) else str(guide)
    return (
        f"No narrative guide found for template {template_id!r}. "
        "Run knowledge generation to create agent_index.yaml."
    )


@tool
def validate_slide_spec(template_id: str, spec: dict) -> dict:
    """Validate a slide specification dict against the element's field schema.

    Returns {valid: bool, errors: [...], warnings: [...]}"""
    element = spec.get("element")
    if not element:
        return {"valid": False, "errors": ["spec.element is required"], "warnings": []}

    path = _template_dir(template_id) / "elements" / f"{element}.yaml"
    if not path.exists():
        return {
            "valid": False,
            "errors": [f"Unknown element {element!r} for template {template_id!r}"],
            "warnings": [],
        }

    data = _load_yaml(path)
    fields = data.get("fields", {})
    if not fields:
        return {
            "valid": False,
            "errors": [f"No fields schema defined for element {element!r}"],
            "warnings": [],
        }
    return _validate_spec_against_schema(spec, fields)


@tool
def build_presentation(
    template_id: str,
    slides: list[dict],
    output_name: str = "output.pptx",
    master_alias: str | None = None,
) -> str:
    """Build a PPTX presentation from a list of slide specification dicts.

    Each dict must have an 'element' key and fields matching that element's schema.
    Validates every slide spec BEFORE building; returns errors if any fail.
    Optional master_alias selects a skeleton theme variant (e.g. "black", "charcoal")."""
    if not slides:
        return "Error: slides list is empty"

    # Pre-validate all specs
    all_errors: list[str] = []
    all_warnings: list[str] = []
    for i, spec in enumerate(slides):
        element = spec.get("element")
        if not element:
            all_errors.append(f"slide[{i}]: missing 'element' key")
            continue
        elem_path = _template_dir(template_id) / "elements" / f"{element}.yaml"
        if not elem_path.exists():
            all_warnings.append(f"slide[{i}]: no schema for element '{element}' — skipping validation")
            continue
        data = _load_yaml(elem_path)
        fields = data.get("fields", {})
        if fields:
            result = _validate_spec_against_schema(spec, fields)
            for e in result.get("errors", []):
                all_errors.append(f"slide[{i}] ({element}): {e}")
            for w in result.get("warnings", []):
                all_warnings.append(f"slide[{i}] ({element}): {w}")

    if all_errors:
        all_warnings.extend(f"(was error) {e}" for e in all_errors)

    assets_root = SKILLS_ROOT / "assets"
    output_path = _PROJECT_ROOT / "decks" / output_name

    try:
        build_deck = _get_build_deck()
        result = build_deck(
            slides,
            assets_root,
            output_path,
            skills_root=SKILLS_ROOT,
            template_id=template_id,
            master_alias=master_alias,
        )
        msg = str(result)
        if all_warnings:
            msg += "\nWarnings:\n" + "\n".join(f"  - {w}" for w in all_warnings)
        return msg
    except Exception as exc:
        return f"Error building presentation: {exc}"


@tool
def read_pdf(file_path: str, max_chars: int = 15000) -> str:
    """Extract text from a PDF file. Returns plain text, truncated to max_chars."""
    try:
        import fitz

        path = Path(file_path)
        if not path.exists():
            return f"Error: file not found: {file_path}"

        parts: list[str] = []
        with fitz.open(path) as doc:
            for page_num, page in enumerate(doc, start=1):
                text = page.get_text().strip()
                if text:
                    parts.append(f"--- Page {page_num} ---\n{text}")

        combined = "\n\n".join(parts)
        if len(combined) > max_chars:
            return combined[:max_chars] + f"\n\n[Truncated at {max_chars} characters]"
        return combined or "[No text extracted from PDF]"
    except Exception as exc:
        return f"Error reading PDF {file_path}: {exc}"


@tool
def read_docx(file_path: str, max_chars: int = 15000) -> str:
    """Extract text from a DOCX file. Returns plain text, truncated to max_chars."""
    try:
        from docx import Document

        path = Path(file_path)
        if not path.exists():
            return f"Error: file not found: {file_path}"

        parts: list[str] = []
        doc = Document(path)
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                parts.append(text)
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))

        combined = "\n\n".join(parts)
        if len(combined) > max_chars:
            return combined[:max_chars] + f"\n\n[Truncated at {max_chars} characters]"
        return combined or "[No text extracted from DOCX]"
    except Exception as exc:
        return f"Error reading DOCX {file_path}: {exc}"


@tool
def read_pptx(file_path: str, max_chars: int = 15000) -> str:
    """Extract slide structure and text from a PPTX file. Returns structured text per slide."""
    try:
        from pptx import Presentation

        path = Path(file_path)
        if not path.exists():
            return f"Error: file not found: {file_path}"

        prs = Presentation(path)
        parts: list[str] = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            layout_name = slide.slide_layout.name if slide.slide_layout else "Unknown"
            texts: list[str] = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
            slide_text = "\n".join(texts) if texts else "[No text on slide]"
            parts.append(f"--- Slide {slide_num} ({layout_name}) ---\n{slide_text}")

        combined = "\n\n".join(parts)
        if len(combined) > max_chars:
            return combined[:max_chars] + f"\n\n[Truncated at {max_chars} characters]"
        return combined or "[No slides found in PPTX]"
    except Exception as exc:
        return f"Error reading PPTX {file_path}: {exc}"


ALL_TOOLS = [
    list_templates,
    get_element_schema,
    list_elements,
    search_icons,
    search_images,
    get_narrative_guide,
    validate_slide_spec,
    build_presentation,
    read_pdf,
    read_docx,
    read_pptx,
]
