"""Extract stock/hero images from PPTX templates into a searchable image catalog."""

from __future__ import annotations

import hashlib
import io
import re
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

import yaml
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

EMU_PER_INCH = 914400

_IMAGE_CATEGORIES = (
    "business",
    "technology",
    "people",
    "abstract",
    "nature",
    "urban",
    "general",
)

_CATEGORY_KEYWORDS: dict[str, list[str]] = {
    "business": ["business", "office", "corporate", "finance", "bank", "sales", "enablement"],
    "technology": ["tech", "technology", "cloud", "data", "digital", "openshift", "sovereignty"],
    "people": ["people", "team", "speaker", "portrait", "headshot"],
    "abstract": ["abstract", "pattern", "wave", "gradient", "background"],
    "nature": ["nature", "landscape", "outdoor", "green"],
    "urban": ["city", "urban", "skyline", "building", "street"],
    "general": [],
}


def _slugify(text: str, max_len: int = 48) -> str:
    slug = re.sub(r"[^a-z0-9]+", "-", text.lower()).strip("-")
    return slug[:max_len] or "stock-image"


def _image_dims(blob: bytes) -> tuple[int, int]:
    try:
        from PIL import Image

        img = Image.open(io.BytesIO(blob))
        return img.size
    except Exception:
        return 0, 0


def _is_stock_image(width_emu: int, height_emu: int) -> bool:
    w_in = width_emu / EMU_PER_INCH
    h_in = height_emu / EMU_PER_INCH
    area = w_in * h_in
    return (w_in > 2.0 and h_in > 1.5) or area > 3.0


def _guess_category(*texts: str) -> str:
    combined = " ".join(t.lower() for t in texts if t)
    for category, keywords in _CATEGORY_KEYWORDS.items():
        if category == "general":
            continue
        if any(kw in combined for kw in keywords):
            return category
    return "general"


def _stub_tags(image_id: str, category: str, layout_name: str) -> list[str]:
    tags = [t for t in image_id.split("-") if len(t) > 2][:5]
    if category not in tags:
        tags.insert(0, category)
    for token in re.findall(r"[a-z]{3,}", layout_name.lower()):
        if token not in tags and token not in {"slide", "layout", "google", "shape"}:
            tags.append(token)
        if len(tags) >= 5:
            break
    return tags[:5]


def _layout_display_name(layout) -> str:
    try:
        return layout.name or "unknown-layout"
    except Exception:
        return "unknown-layout"


def _extract_from_pptx(
    pptx_path: Path,
    images_dir: Path,
    seen_hashes: set[str],
    catalog_entries: list[dict[str, Any]],
    stock_index: int,
) -> int:
    """Extract stock images from one PPTX; return next stock_index."""
    prs = Presentation(str(pptx_path))

    for layout in prs.slide_layouts:
        layout_name = _layout_display_name(layout)
        for shape in layout.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            if not _is_stock_image(shape.width, shape.height):
                continue

            try:
                blob = shape.image.blob
            except Exception:
                continue

            digest = hashlib.sha1(blob).hexdigest()[:12]
            if digest in seen_hashes:
                continue
            seen_hashes.add(digest)

            shape_name = getattr(shape, "name", "") or f"stock-{stock_index}"
            base_id = _slugify(shape_name)
            if base_id.startswith("google-shape"):
                base_id = f"stock-{stock_index}"
            image_id = base_id
            suffix = 2
            existing_ids = {entry["id"] for entry in catalog_entries}
            while image_id in existing_ids:
                image_id = f"{base_id}-{suffix}"
                suffix += 1

            ext = (getattr(shape.image, "ext", None) or "png").lower()
            if ext == "jpeg":
                ext = "jpg"
            dest = images_dir / f"{image_id}.{ext}"
            dest.write_bytes(blob)

            width_px, height_px = _image_dims(blob)
            category = _guess_category(image_id, layout_name, shape_name)
            tags = _stub_tags(image_id, category, layout_name)

            catalog_entries.append({
                "id": image_id,
                "file": f"images/{dest.name}",
                "width_px": width_px,
                "height_px": height_px,
                "source_layout": layout_name,
                "source_file": pptx_path.name,
                "tags": tags,
                "category": category,
                "description": (
                    f"Stock image from {layout_name}, "
                    f"suitable for title-block or image-content slides"
                ),
            })
            stock_index += 1

    return stock_index


def extract_stock_images(
    pptx_paths: list[Path | str],
    output_root: Path,
) -> Path:
    """Extract stock images from PPTX files and write image_catalog.yaml."""
    assets_dir = output_root / "assets"
    images_dir = assets_dir / "images"
    images_dir.mkdir(parents=True, exist_ok=True)

    catalog_entries: list[dict[str, Any]] = []
    seen_hashes: set[str] = set()
    stock_index = 1

    for raw_path in pptx_paths:
        path = Path(raw_path)
        if not path.exists() or path.suffix.lower() != ".pptx":
            continue
        try:
            stock_index = _extract_from_pptx(
                path, images_dir, seen_hashes, catalog_entries, stock_index,
            )
        except Exception as exc:
            print(f"  WARNING: stock image extraction failed for {path}: {exc}")

    catalog_path = assets_dir / "image_catalog.yaml"
    data = {
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "source_files": [str(Path(p).name) for p in pptx_paths if Path(p).exists()],
        "images": catalog_entries,
    }
    with catalog_path.open("w", encoding="utf-8") as fh:
        yaml.dump(data, fh, allow_unicode=True, sort_keys=False, width=120)

    return catalog_path


def discover_template_pptx(project_root: Path | None = None) -> list[Path]:
    """Find sovereignty and standard template PPTX files under templates/."""
    root = project_root or Path(__file__).resolve().parent.parent.parent
    templates_dir = root / "templates"
    if not templates_dir.is_dir():
        return []
    return sorted(templates_dir.glob("*.pptx"))


def extract_stock_images_for_project(
    primary_pptx: Path | str | None,
    output_root: Path,
    *,
    include_templates: bool = True,
) -> Path:
    """Extract from a primary PPTX plus optional templates/ directory scan."""
    project_root = Path(__file__).resolve().parent.parent.parent
    paths: list[Path] = []
    if primary_pptx:
        p = Path(primary_pptx)
        if p.exists():
            paths.append(p)
    if include_templates:
        for tpl in discover_template_pptx(project_root):
            if tpl not in paths:
                paths.append(tpl)
    if not paths:
        paths = discover_template_pptx(project_root)
    return extract_stock_images(paths, output_root)
