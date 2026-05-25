"""Extract complete visual blueprints from every slide layout in a PPTX.

For each layout we record:
  - Every decorative shape (non-placeholder) with:
      * kind: picture | line | rect | freeform | textbox
      * position / size as percentage of slide dimensions
      * fill: solid color hex | transparent | image_ref
      * line: color hex, width_pt, dash_style
      * image_ref: sha1 filename if it is a picture
  - Every placeholder with its idx, type, and geometry

The output is a dict keyed by layout name that can be serialised to YAML and
embedded in SKILL.md, making the skills fully self-contained.
"""

from __future__ import annotations

import hashlib
import re
import zipfile
from pathlib import Path
from typing import Any

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt


NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}

# Dash styles the OOXML uses
DASH_MAP = {
    "solid": "solid",
    "dash": "dash",
    "dashDot": "dash-dot",
    "dot": "dot",
    "lgDash": "long-dash",
    "lgDashDot": "long-dash-dot",
    "sysDash": "sys-dash",
    "sysDot": "sys-dot",
}


def _hex(rgb_val) -> str:
    try:
        return f"#{rgb_val}"
    except Exception:
        return "none"


def _dash_style(ln_el) -> str:
    if ln_el is None:
        return "solid"
    dash_el = ln_el.find("a:prstDash", NS)
    if dash_el is not None:
        return DASH_MAP.get(dash_el.get("val", "solid"), "solid")
    cust = ln_el.find("a:custDash", NS)
    if cust is not None:
        return "custom-dash"
    return "solid"


def _line_info(shape) -> dict:
    try:
        line = shape.line
        info: dict = {}
        try:
            info["color"] = _hex(line.color.rgb)
        except Exception:
            info["color"] = "none"
        info["width_pt"] = round(line.width.pt, 2) if line.width else 0
        ln_el = shape._element.find(".//a:ln", NS)
        info["dash"] = _dash_style(ln_el)
        return info
    except Exception:
        return {"color": "none", "width_pt": 0, "dash": "solid"}


def _fill_info(shape) -> dict:
    try:
        ft = shape.fill.type
        if ft is None:
            return {"type": "none"}
        name = ft.name
        if name == "SOLID":
            try:
                return {"type": "solid", "color": _hex(shape.fill.fore_color.rgb)}
            except Exception:
                return {"type": "solid", "color": "theme"}
        if name == "BACKGROUND":
            return {"type": "transparent"}
        return {"type": name.lower()}
    except Exception:
        return {"type": "unknown"}


def _placeholder_font_meta(sp) -> dict:
    """Extract default typography from a placeholder shape's txBody."""
    meta: dict = {}

    for xpath in [
        ".//p:txBody/a:lstStyle/a:lvl1pPr/a:defRPr",
        ".//p:txBody/a:p/a:pPr/a:defRPr",
        ".//p:txBody/a:p/a:r/a:rPr",
    ]:
        rpr = sp.find(xpath, NS) if hasattr(sp, "find") else None
        if rpr is None:
            continue

        sz = rpr.get("sz")
        if sz:
            meta["font_size_pt"] = int(sz) / 100

        b = rpr.get("b")
        if b:
            meta["bold"] = b == "1" or b.lower() == "true"

        typeface_el = rpr.find("a:latin", NS)
        if typeface_el is not None:
            tf = typeface_el.get("typeface")
            if tf and not tf.startswith("+"):
                meta["font_name"] = tf

        fill = rpr.find("a:solidFill", NS)
        if fill is not None:
            srgb = fill.find("a:srgbClr", NS)
            scheme = fill.find("a:schemeClr", NS)
            if srgb is not None:
                meta["font_color"] = f"#{srgb.get('val', '000000')}"
            elif scheme is not None:
                meta["font_color"] = "theme"
                meta["color_scheme"] = scheme.get("val")

        if meta:
            break

    for xpath in [
        ".//p:txBody/a:lstStyle/a:lvl1pPr",
        ".//p:txBody/a:p/a:pPr",
    ]:
        ppr = sp.find(xpath, NS) if hasattr(sp, "find") else None
        if ppr is not None:
            algn = ppr.get("algn")
            if algn:
                meta["alignment"] = algn
            break

    return meta


def _geom(shape, W: int, H: int) -> dict:
    def pct(v, dim):
        return round(int(v) / dim * 100, 2) if v is not None else None

    return {
        "x_pct": pct(shape.left, W),
        "y_pct": pct(shape.top, H),
        "w_pct": pct(shape.width, W),
        "h_pct": pct(shape.height, H),
        # absolute EMU as plain int (no pptx.util.Emu objects)
        "x_emu": int(shape.left)   if shape.left   is not None else None,
        "y_emu": int(shape.top)    if shape.top     is not None else None,
        "w_emu": int(shape.width)  if shape.width   is not None else None,
        "h_emu": int(shape.height) if shape.height  is not None else None,
    }


def _extract_layout_media(pptx_path: str, layout_idx: int) -> dict[str, bytes]:
    """Return {rId: image_bytes} for a specific layout's part."""
    z = zipfile.ZipFile(pptx_path)
    prs = Presentation(pptx_path)
    layout = prs.slide_layouts[layout_idx]
    media: dict[str, bytes] = {}
    for rId, rel in layout.part.rels.items():
        if "image" in rel.reltype.lower():
            try:
                blob = rel.target_part.blob
                media[rId] = blob
            except Exception:
                pass
    return media


def _extract_bg_image(part, img_dir: Path) -> dict | None:
    """Extract the background image from a layout/master part's <p:bg> element.

    Returns a dict with image_ref and fill_mode, or None if no image background.
    """
    bg_el = part._element.find(
        ".//{http://schemas.openxmlformats.org/presentationml/2006/main}bg"
    )
    if bg_el is None:
        return None

    # Look for blipFill (image background)
    blip = bg_el.find(".//a:blip", NS)
    if blip is None:
        # Check for solid fill colour
        solidFill = bg_el.find(".//a:solidFill", NS)
        if solidFill is not None:
            srgb = solidFill.find("a:srgbClr", NS)
            if srgb is not None:
                return {"type": "solid", "color": f"#{srgb.get('val', '000000')}"}
        return None

    rid = blip.get(f"{{{NS['r']}}}embed")
    if not rid or rid not in part.rels:
        return None

    try:
        blob = part.rels[rid].target_part.blob
    except Exception:
        return None

    sha = hashlib.sha1(blob).hexdigest()[:12]
    fname = f"bg_{sha}.png"
    dest = img_dir / fname
    if not dest.exists():
        dest.write_bytes(blob)

    # Detect fill mode (stretch vs tile)
    stretch = bg_el.find(".//a:stretch", NS)
    tile    = bg_el.find(".//a:tile", NS)
    mode = "stretch" if stretch is not None else ("tile" if tile is not None else "stretch")

    return {
        "type":       "image",
        "image_ref":  fname,
        "fill_mode":  mode,
        "size_bytes": len(blob),
    }


_GENERIC_NAMES = re.compile(
    r"^(TITLE|CUSTOM|BLANK|LAYOUT|SLIDE|CONTENT|SECTION|MASTER)[\d_]*$",
    re.IGNORECASE,
)

# Placeholder type idx → human label
_PH_IDX_LABEL = {
    0: "Title",
    1: "Content",
    2: "Body",
    3: "Center Title",
    10: "Date",
    11: "Footer",
    12: "Slide Number",
    13: "Subtitle",
    14: "Header",
    15: "Media Clip",
}

_PH_TYPE_LABEL = {
    "TITLE": "Title",
    "CENTER_TITLE": "Center Title",
    "BODY": "Body",
    "SUBTITLE": "Subtitle",
    "PICTURE": "Picture",
    "CHART": "Chart",
    "TABLE": "Table",
    "MEDIA_CLIP": "Media Clip",
    "OBJECT": "Object",
}


def _infer_display_name(internal_name: str, placeholders: list[dict]) -> str:
    """Return a human-readable layout name.

    If the internal name is already descriptive (e.g. 'Title block_1',
    'Blank page_1') keep it as-is.  When it looks generic (TITLE, CUSTOM_2,
    BLANK_3) derive a name from the placeholder types present.
    """
    base = internal_name.split(":")[-1] if ":" in internal_name else internal_name
    if not _GENERIC_NAMES.match(base):
        return internal_name  # already readable

    ph_labels = []
    for ph in placeholders:
        label = _PH_TYPE_LABEL.get(ph.get("type", ""), None)
        if label is None:
            label = _PH_IDX_LABEL.get(ph.get("idx", -1), None)
        if label and label not in ph_labels:
            ph_labels.append(label)

    # Strip decoration placeholders (date/footer/page-number) for the name
    content_labels = [l for l in ph_labels if l not in ("Date", "Footer", "Slide Number", "Header")]
    if content_labels:
        return " + ".join(content_labels) + " Layout"
    if ph_labels:
        return " + ".join(ph_labels) + " Layout"
    return internal_name  # nothing useful found, keep original


def extract_blueprints(pptx_path: str, assets_dir: Path) -> dict[str, Any]:
    """
    Extract layout blueprints from all layouts in the PPTX.

    Saves image blobs to assets_dir/layout-images/ and returns structured dict.
    """
    img_dir = assets_dir / "layout-images"
    img_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(pptx_path)
    W = int(prs.slide_width)
    H = int(prs.slide_height)
    z = zipfile.ZipFile(pptx_path)
    all_media = {f: z.read(f) for f in z.namelist() if f.startswith("ppt/media/")}

    blueprints: dict[str, Any] = {}

    # ── Extract master-level decorative shapes (red stripe, logo, etc.) ──
    for mi, master in enumerate(prs.slide_masters):
        master_shapes: list[dict] = []
        rId_to_blob_m: dict[str, bytes] = {}
        for rId, rel in master.part.rels.items():
            if "image" in rel.reltype.lower():
                try:
                    rId_to_blob_m[rId] = rel.target_part.blob
                except Exception:
                    pass

        for shape in master.shapes:
            if shape.is_placeholder:
                continue
            st = shape.shape_type
            entry: dict[str, Any] = {
                "name": shape.name,
                "kind": "unknown",
                **_geom(shape, W, H),
            }
            if st == MSO_SHAPE_TYPE.PICTURE:
                entry["kind"] = "picture"
                blip = shape._element.find(".//a:blip", NS)
                rid = blip.get(f"{{{NS['r']}}}embed") if blip is not None else None
                blob = rId_to_blob_m.get(rid) if rid else None
                if blob:
                    sha = hashlib.sha1(blob).hexdigest()[:12]
                    fname = f"{sha}.png"
                    (img_dir / fname).write_bytes(blob)
                    entry["image_ref"] = fname
                    entry["image_size_bytes"] = len(blob)
                else:
                    entry["image_ref"] = None
            elif st == MSO_SHAPE_TYPE.LINE:
                entry["kind"] = "line"
                entry["line"] = _line_info(shape)
            elif st in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM):
                entry["kind"] = "shape"
                entry["fill"] = _fill_info(shape)
                entry["line"] = _line_info(shape)
                prstGeom = shape._element.find(".//a:prstGeom", NS)
                if prstGeom is not None:
                    entry["preset_geometry"] = prstGeom.get("prst")
            elif st == MSO_SHAPE_TYPE.TEXT_BOX:
                entry["kind"] = "textbox"
                entry["fill"] = _fill_info(shape)
            else:
                entry["kind"] = str(st)
                entry["fill"] = _fill_info(shape)
                entry["line"] = _line_info(shape)
            master_shapes.append(entry)

        master_bg = _extract_bg_image(master.part, img_dir)
        key = f"_MASTER_{mi}"
        blueprints[key] = {
            "index":             0,
            "master":            mi,
            "name":              f"Slide Master {mi}",
            "display_name":      f"Slide Master {mi}",
            "key":               key,
            "slide_width_emu":   int(W),
            "slide_height_emu":  int(H),
            "background":        master_bg,
            "decorative_shapes": master_shapes,
            "placeholders":      [],
        }

    # ── Iterate ALL slide layouts ──
    layout_counter = 0
    all_layouts: list[tuple[int, int, Any]] = []
    for mi, master in enumerate(prs.slide_masters):
        for li, layout in enumerate(master.slide_layouts):
            all_layouts.append((mi, li, layout))

    for (mi, li, layout) in all_layouts:
        layout_counter += 1
        shapes_out = []
        placeholders_out = []

        # Build rId → media blob map for this layout
        rId_to_blob: dict[str, bytes] = {}
        for rId, rel in layout.part.rels.items():
            if "image" in rel.reltype.lower():
                try:
                    blob = rel.target_part.blob
                    rId_to_blob[rId] = blob
                except Exception:
                    pass

        for shape in layout.shapes:
            if shape.is_placeholder:
                ph = shape.placeholder_format
                ph_dict = {
                    "idx":  ph.idx,
                    "type": ph.type.name,
                    "name": shape.name,
                    **_geom(shape, W, H),
                }
                ph_dict.update(_placeholder_font_meta(shape._element))
                placeholders_out.append(ph_dict)
                continue

            st = shape.shape_type
            entry: dict[str, Any] = {
                "name": shape.name,
                "kind": "unknown",
                **_geom(shape, W, H),
            }

            if st == MSO_SHAPE_TYPE.PICTURE:
                entry["kind"] = "picture"
                # resolve rId from XML
                blip = shape._element.find(".//a:blip", NS)
                rid = blip.get(f"{{{NS['r']}}}embed") if blip is not None else None
                blob = rId_to_blob.get(rid) if rid else None
                if blob:
                    sha = hashlib.sha1(blob).hexdigest()[:12]
                    ext = "png"  # all are PNG in this template
                    fname = f"{sha}.{ext}"
                    (img_dir / fname).write_bytes(blob)
                    entry["image_ref"] = fname
                    entry["image_size_bytes"] = len(blob)
                else:
                    entry["image_ref"] = None

            elif st == MSO_SHAPE_TYPE.LINE:
                entry["kind"] = "line"
                entry["line"] = _line_info(shape)

            elif st == MSO_SHAPE_TYPE.TEXT_BOX:
                entry["kind"] = "textbox"
                entry["fill"] = _fill_info(shape)
                try:
                    tf = shape.text_frame
                    para = tf.paragraphs[0] if tf.paragraphs else None
                    if para and para.runs:
                        run = para.runs[0]
                        entry["text_sample"] = run.text[:80]
                        entry["font_name"]   = run.font.name
                        entry["font_size_pt"] = round(run.font.size.pt, 1) if run.font.size else None
                        try:
                            entry["font_color"] = _hex(run.font.color.rgb)
                        except Exception:
                            entry["font_color"] = "theme"
                except Exception:
                    pass

            elif st in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM):
                entry["kind"] = "shape"
                entry["fill"] = _fill_info(shape)
                entry["line"] = _line_info(shape)
                # Try to get preset geometry name
                prstGeom = shape._element.find(".//a:prstGeom", NS)
                if prstGeom is not None:
                    entry["preset_geometry"] = prstGeom.get("prst")

            else:
                entry["kind"] = str(st)
                entry["fill"] = _fill_info(shape)
                entry["line"] = _line_info(shape)

            shapes_out.append(entry)

        # Extract background image (the wave/dots/dashes live here)
        bg_info = _extract_bg_image(layout.part, img_dir)

        # Use master_index.layout_name as key to avoid collisions across masters
        key = f"M{mi}:{layout.name}" if mi > 0 else layout.name
        display_name = _infer_display_name(key, placeholders_out)

        blueprints[key] = {
            "index":        layout_counter,
            "master":       mi,
            "name":         layout.name,
            "display_name": display_name,
            "key":          key,
            "slide_width_emu":  int(W),
            "slide_height_emu": int(H),
            "background":        bg_info,
            "decorative_shapes": shapes_out,
            "placeholders":      placeholders_out,
        }

    return blueprints


def save_blueprints(pptx_path: str, output_root: Path, template_id: str) -> Path:
    """
    Extract blueprints and write them to:
      output_root/templates/<template_id>/layouts/blueprints.yaml
    Also saves extracted images to:
      output_root/assets/layout-images/
    """
    import yaml

    assets_dir = output_root / "assets"
    blueprints = extract_blueprints(pptx_path, assets_dir)

    out_dir = output_root / "templates" / template_id / "layouts"
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / "blueprints.yaml"

    with open(out_path, "w") as f:
        yaml.dump({"layouts": blueprints}, f, allow_unicode=True, sort_keys=False)

    n_layouts = len(blueprints)
    n_shapes  = sum(len(v["decorative_shapes"]) for v in blueprints.values())
    n_images  = sum(
        1 for v in blueprints.values()
        for s in v["decorative_shapes"] if s.get("kind") == "picture"
    )
    print(f"  Blueprints: {n_layouts} layouts, {n_shapes} decorative shapes, {n_images} images")
    print(f"  Saved: {out_path}")
    return out_path
