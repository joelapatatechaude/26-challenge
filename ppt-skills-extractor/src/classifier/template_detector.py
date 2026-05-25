"""
Fingerprint a PPTX file and match it to a known template in templates_registry.yaml.
If no match is found, a draft entry is created and the user is prompted to name it.
"""

from __future__ import annotations

import hashlib
import re
import uuid
from pathlib import Path
from typing import Optional

import yaml
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from rich.console import Console
from rich.prompt import Prompt

from schemas import TemplateProfile

console = Console()

REGISTRY_PATH = Path(__file__).parents[2] / "config" / "templates_registry.yaml"


# ---------------------------------------------------------------------------
# Fingerprinting helpers
# ---------------------------------------------------------------------------

def _theme_colors(prs: Presentation) -> dict[str, str]:
    """Extract named theme colours from the slide master's theme XML."""
    colors: dict[str, str] = {}
    try:
        from lxml import etree
        theme_el = prs.slide_master.theme_color_map._element
        # Walk up to find the theme element
        root = theme_el
        while root.getparent() is not None:
            root = root.getparent()
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        for srgb in root.findall(".//a:srgbClr", ns):
            val = srgb.get("val", "")
            parent = srgb.getparent()
            if parent is not None:
                tag = re.sub(r"\{.*?\}", "", parent.tag)
                if val:
                    colors[tag] = f"#{val.upper()}"
    except Exception:
        pass
    return colors


def _theme_fonts(prs: Presentation) -> list[str]:
    """Extract font names from the slide master theme."""
    fonts: list[str] = []
    try:
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        root = prs.slide_master._element
        for latin in root.findall(".//a:latin", ns):
            typeface = latin.get("typeface", "")
            if typeface and not typeface.startswith("+") and typeface not in fonts:
                fonts.append(typeface)
    except Exception:
        pass
    return fonts[:5]


def _fonts_from_placeholders(prs: Presentation) -> tuple[str | None, str | None]:
    """Scan slide layout placeholders to determine heading and body fonts.

    Title placeholders (idx=0) → heading font.
    Body placeholders (idx=1) → body font.
    Prefers explicit font names over theme references (+mj-lt etc.).
    """
    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main",
          "p": "http://schemas.openxmlformats.org/presentationml/2006/main"}

    heading_font: str | None = None
    body_font: str | None = None

    def _pick_font(sp_el) -> str | None:
        for xpath in [
            ".//p:txBody/a:lstStyle/a:lvl1pPr/a:defRPr/a:latin",
            ".//p:txBody/a:p/a:r/a:rPr/a:latin",
        ]:
            el = sp_el.find(xpath, ns)
            if el is not None:
                tf = el.get("typeface", "")
                if tf and not tf.startswith("+"):
                    return tf
        return None

    sources = list(prs.slide_masters) + [
        layout
        for master in prs.slide_masters
        for layout in master.slide_layouts
    ]

    for part in sources:
        for shape in part.shapes:
            if not shape.is_placeholder:
                continue
            ph = shape.placeholder_format
            font = _pick_font(shape._element)
            if font:
                if ph.idx == 0 and heading_font is None:
                    heading_font = font
                elif ph.idx in (1, 2) and body_font is None:
                    body_font = font
        if heading_font and body_font:
            break

    return (heading_font, body_font)


def _theme_font_pair(prs: Presentation) -> tuple[str | None, str | None]:
    """Extract majorFont and minorFont typefaces from theme XML."""
    from lxml import etree

    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    try:
        master = prs.slide_masters[0]
        for rel in master.part.rels.values():
            if "theme" in rel.reltype:
                root = etree.fromstring(rel.target_part.blob)
                major = root.findall(".//a:majorFont/a:latin", ns)
                minor = root.findall(".//a:minorFont/a:latin", ns)
                mj = major[0].get("typeface") if major else None
                mn = minor[0].get("typeface") if minor else None
                if mj and mj.startswith("+"):
                    mj = None
                if mn and mn.startswith("+"):
                    mn = None
                return (mj, mn)
    except Exception:
        pass
    return (None, None)


def _logo_hash(prs: Presentation) -> Optional[str]:
    """Hash the first picture shape found on the slide master (likely the logo)."""
    try:
        for shape in prs.slide_master.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                return hashlib.sha1(shape.image.blob).hexdigest()[:16]
    except Exception:
        pass
    return None


def _primary_color(prs: Presentation) -> str:
    """Best-guess the primary brand colour (accent1 or dk1)."""
    colors = _theme_colors(prs)
    for key in ("accent1", "dk1", "lt2"):
        if key in colors:
            return colors[key]
    # Fallback: look for red-ish fill on master shapes
    try:
        for shape in prs.slide_master.shapes:
            fc = None
            try:
                fc = shape.fill.fore_color.rgb
            except Exception:
                pass
            if fc:
                hex_c = f"#{fc}".upper()
                r = int(str(fc)[:2], 16)
                g = int(str(fc)[2:4], 16)
                b = int(str(fc)[4:6], 16)
                if r > 150 and g < 50 and b < 50:
                    return hex_c
    except Exception:
        pass
    return "#EE0000"


def _background_color(prs: Presentation) -> Optional[str]:
    """Detect a dominant dark background if present (for dark templates)."""
    try:
        bg = prs.slide_master.background.fill
        if bg.fore_color and bg.fore_color.rgb:
            rgb = str(bg.fore_color.rgb)
            r = int(rgb[:2], 16)
            g = int(rgb[2:4], 16)
            b = int(rgb[4:6], 16)
            if r < 50 and g < 50 and b < 50:
                return f"#{rgb.upper()}"
    except Exception:
        pass
    return None


def _accent_colors(prs: Presentation) -> list[str]:
    """Return accent1-accent6 colors for discriminating similar templates."""
    colors = _theme_colors(prs)
    return [colors.get(f"accent{i}", "") for i in range(1, 7)]


# Fonts that are theme defaults and frequently overridden by actual slide fonts.
_GENERIC_FONTS = {"arial", "calibri", "times new roman", "helvetica", "tahoma", "verdana"}


def _is_generic_font(name: str | None) -> bool:
    return not name or name.strip().lower() in _GENERIC_FONTS


def _fingerprint(prs: Presentation) -> dict:
    mj, mn = _theme_font_pair(prs)

    # Placeholder-level fonts reflect what's actually rendered on slides,
    # so always prefer them over generic theme defaults (Arial, Calibri, etc.)
    # that are often just placeholders in the theme XML.
    ph_heading, ph_body = _fonts_from_placeholders(prs)
    if ph_heading and _is_generic_font(mj):
        mj = ph_heading
    if ph_body and _is_generic_font(mn):
        mn = ph_body

    return {
        "primary_color": _primary_color(prs),
        "background_dark": _background_color(prs),
        "accent_colors": _accent_colors(prs),
        "fonts": _theme_fonts(prs),
        "font_heading": mj,
        "font_body": mn,
        "logo_hash": _logo_hash(prs),
        "slide_width_emu": prs.slide_width,
        "slide_height_emu": prs.slide_height,
    }


# ---------------------------------------------------------------------------
# Registry helpers
# ---------------------------------------------------------------------------

def _load_registry() -> dict:
    if REGISTRY_PATH.exists():
        data = yaml.safe_load(REGISTRY_PATH.read_text()) or {}
        return data
    return {"templates": []}


def _save_registry(data: dict) -> None:
    REGISTRY_PATH.write_text(yaml.dump(data, default_flow_style=False, allow_unicode=True))


def _slug_from_filename(path: str | Path) -> str:
    """Derive a template ID slug from a PPTX filename."""
    name = Path(path).stem
    for suffix in [" presentation template", " template", " presentation"]:
        name = name.replace(suffix, "").replace(suffix.title(), "")
    slug = re.sub(r"[^a-z0-9]+", "-", name.lower()).strip("-")
    slug = re.sub(r"-+", "-", slug)
    return slug


def _has_accent_colors(accents: list[str] | None) -> bool:
    """True when the list contains at least one non-empty accent colour."""
    return bool(accents and any(c for c in accents if c))


def _accent_overlap(fp_accents: list[str], entry_accents: list[str]) -> int:
    """Count how many accent colors overlap between fingerprint and registry entry."""
    a = {c.upper() for c in fp_accents if c}
    b = {c.upper() for c in entry_accents if c}
    return len(a & b)


def _match(fp: dict, entry: dict) -> bool:
    """Match a fingerprint against a registry entry.

    Priority order:
      1. Logo hash (exact)
      2. Primary colour + ≥2 accent colours
      3. Primary colour + slide dimensions + font overlap (handles templates
         where accent colours are not exposed in the theme XML)
      4. Legacy: primary + fonts, no accent_colors in the entry
    """
    if fp.get("logo_hash") and fp["logo_hash"] == entry.get("logo_hash"):
        return True

    color_match = fp["primary_color"].upper() == entry.get("primary_color", "").upper()
    fp_accents = [c for c in fp.get("accent_colors", []) if c]
    accent_score = _accent_overlap(fp.get("accent_colors", []), entry.get("accent_colors", []))

    # Strong: ≥2 accent colours in common
    if color_match and accent_score >= 2:
        return True

    # Dimension-based: slide size is reliable when both templates share red+no-accents
    dim_match = (
        fp.get("slide_width_emu") == entry.get("slide_width_emu")
        and fp.get("slide_height_emu") == entry.get("slide_height_emu")
    )
    font_overlap = bool(set(fp.get("fonts", [])) & set(entry.get("fonts", [])))

    if color_match and dim_match and font_overlap and not _has_accent_colors(fp.get("accent_colors")) and not _has_accent_colors(entry.get("accent_colors")):
        return True

    # Legacy fallback: primary + fonts, no accent_colors registered
    if color_match and font_overlap and not _has_accent_colors(entry.get("accent_colors")):
        return True

    return False


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def detect_template(pptx_path: Path, interactive: bool = False) -> TemplateProfile:
    """
    Fingerprint the PPTX and return a TemplateProfile.
    If unrecognised and interactive=True, prompt the user to name the new template.
    """
    prs = Presentation(str(pptx_path))
    fp = _fingerprint(prs)
    registry = _load_registry()
    templates = registry.get("templates") or []

    for entry in templates:
        if _match(fp, entry):
            console.print(f"  [green]Template matched:[/green] {entry['name']} ({entry['id']})")
            return TemplateProfile(
                template_id=entry["id"],
                name=entry["name"],
                primary_color=entry.get("primary_color", fp["primary_color"]),
                background_dark=entry.get("background_dark"),
                background_light=entry.get("background_light", "#f5f5f5"),
                font_heading=fp.get("font_heading") or (fp["fonts"][0] if fp["fonts"] else "Red Hat Display"),
                font_body=fp.get("font_body") or (fp["fonts"][0] if fp["fonts"] else "Red Hat Display"),
                logo_image_hash=fp.get("logo_hash"),
                slide_width_emu=int(fp["slide_width_emu"]),
                slide_height_emu=int(fp["slide_height_emu"]),
            )

    # No match — create a draft entry
    if interactive:
        template_id = f"template-{uuid.uuid4().hex[:8]}"
        console.print(f"\n  [yellow]Unrecognised template[/yellow] in [bold]{pptx_path.name}[/bold]")
        console.print(f"  Detected: primary={fp['primary_color']}, fonts={fp['fonts']}, logo_hash={fp.get('logo_hash')}")
        name = Prompt.ask("  Enter a short name for this template", default=template_id)
        template_id = re.sub(r"[^a-z0-9]+", "-", name.lower()).strip("-")
    else:
        template_id = _slug_from_filename(pptx_path)
        name = template_id.replace("-", " ").title()

    existing_ids = {t["id"] for t in templates}
    if template_id in existing_ids:
        suffix = 2
        candidate_id = f"{template_id}-{suffix}"
        while candidate_id in existing_ids:
            suffix += 1
            candidate_id = f"{template_id}-{suffix}"
        template_id = candidate_id
        if not interactive:
            name = f"{name} {suffix}"

    if not interactive:
        console.print(f"  [green]Auto-registered new template:[/green] {template_id}")

    entry = {
        "id": template_id,
        "name": name,
        "primary_color": fp["primary_color"],
        "background_dark": fp.get("background_dark"),
        "background_light": "#f5f5f5",
        "accent_colors": fp.get("accent_colors", []),
        "fonts": fp["fonts"],
        "logo_hash": fp.get("logo_hash"),
        "slide_width_emu": int(fp["slide_width_emu"]),
        "slide_height_emu": int(fp["slide_height_emu"]),
    }
    templates.append(entry)
    registry["templates"] = templates
    _save_registry(registry)
    if interactive:
        console.print(f"  [green]New template registered:[/green] {name} ({template_id})")

    return TemplateProfile(
        template_id=template_id,
        name=name,
        primary_color=fp["primary_color"],
        background_dark=fp.get("background_dark"),
        background_light="#f5f5f5",
        font_heading=fp.get("font_heading") or (fp["fonts"][0] if fp["fonts"] else "Red Hat Display"),
        font_body=fp.get("font_body") or (fp["fonts"][0] if fp["fonts"] else "Red Hat Display"),
        logo_image_hash=fp.get("logo_hash"),
        slide_width_emu=int(fp["slide_width_emu"]),
        slide_height_emu=int(fp["slide_height_emu"]),
    )
