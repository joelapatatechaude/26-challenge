"""Build a real PPTX file from a deck spec (parsed markdown or dict).

Each slide is rendered using element-specific drawing routines that follow
the rules in skills-output/templates/red-hat-standard/elements/*.yaml.

Usage:
    from generator.deck_builder import build_deck
    build_deck(spec, skills_root=Path("skills-output"), output_path=Path("out.pptx"))
"""

from __future__ import annotations

import logging
import re
from collections.abc import Callable
from pathlib import Path
from typing import Any

import yaml
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Emu, Inches, Pt

from generator.brand_palette import RHColor
from generator.layout_context import (
    LayoutMetrics,
    Rect,
    SlideRenderContext,
    TypographyScale,
)

log = logging.getLogger("ppt.deck_builder")


def resolve_slide_layout_name(
    prs: Presentation,
    layout_name: str,
    template_id: str | None = None,
    *,
    on_warn: Callable[[str], None] | None = None,
) -> str:
    """Validate layout name against the presentation; fall back to first layout."""
    available_layouts = {layout.name for layout in prs.slide_layouts}
    if layout_name not in available_layouts:
        log.warning(
            "Layout '%s' not found in template '%s'. Available: %s. Falling back to first layout.",
            layout_name,
            template_id,
            sorted(available_layouts),
        )
        if on_warn:
            on_warn(layout_name)
        return prs.slide_layouts[0].name
    return layout_name


# ---------------------------------------------------------------------------
# Template brand — loaded from template.yaml, falls back to RH Standard
# ---------------------------------------------------------------------------

from dataclasses import dataclass, field


@dataclass
class TemplateBrand:
    """All template-specific visual constants derived from template.yaml."""
    # Colours
    primary:     RGBColor = field(default_factory=lambda: RGBColor(0xEE, 0x00, 0x00))
    dark:        RGBColor = field(default_factory=lambda: RGBColor(0x5F, 0x00, 0x00))
    accent:      RGBColor = field(default_factory=lambda: RGBColor(0x37, 0xA3, 0xA3))  # teal-50
    black:       RGBColor = field(default_factory=lambda: RGBColor(0x15, 0x15, 0x15))  # gray-95
    white:       RGBColor = field(default_factory=lambda: RGBColor(0xFF, 0xFF, 0xFF))
    light_bg:    RGBColor = field(default_factory=lambda: RGBColor(0xF2, 0xF2, 0xF2))  # gray-10
    dark_grey:   RGBColor = field(default_factory=lambda: RGBColor(0x4D, 0x4D, 0x4D))  # gray-60
    mid_grey:    RGBColor = field(default_factory=lambda: RGBColor(0x70, 0x70, 0x70))  # gray-50
    amber:       RGBColor = field(default_factory=lambda: RGBColor(0xF5, 0x92, 0x1B))  # orange-40
    # Fonts
    font_heading: str = "Red Hat Display"
    font_body:    str = "Red Hat Text"
    # Flags
    has_blueprint_bg: bool = False   # True → renderers must NOT paint slide background

    @classmethod
    def from_template_yaml(cls, path: Path) -> "TemplateBrand":
        if not path.exists():
            return cls()
        with path.open() as f:
            data = yaml.safe_load(f) or {}
        theme = data.get("theme", {})
        colors = theme.get("colors", {})

        def c(hex_str: str, fallback: RGBColor) -> RGBColor:
            try:
                return _rgb(hex_str.lstrip("#")) if hex_str else fallback
            except Exception:
                return fallback

        primary = c(theme.get("primary_color", ""), RGBColor(0xEE, 0x00, 0x00))
        teal    = c(theme.get("teal_accent") or colors.get("accent5", ""), primary)
        dark    = c(theme.get("dark_accent") or colors.get("dk2", ""),
                    RGBColor(0x5F, 0x00, 0x00))
        light   = c(theme.get("background_light") or colors.get("lt2", ""),
                    RGBColor(0xF5, 0xF5, 0xF5))
        amber   = c(theme.get("amber_accent") or colors.get("accent1", ""),
                    RGBColor(0xFF, 0xAB, 0x40))

        # Always use brand fonts per Red Hat brand standards
        fh = "Red Hat Display"
        fb = "Red Hat Text"

        return cls(
            primary=primary, dark=dark, accent=teal,
            light_bg=light, amber=amber,
            font_heading=fh, font_body=fb,
        )


# Fallback brand (Red Hat Standard)
_DEFAULT_BRAND = TemplateBrand()

# Convenience aliases still referenced by helpers (resolved at render time via brand param)
BLACK    = RGBColor(0x15, 0x15, 0x15)  # gray-95
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
MID_GREY = RGBColor(0x70, 0x70, 0x70)  # gray-50

# Slide canvas — widescreen 16:9
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# Standard margins
MARGIN_L = Inches(0.5)
MARGIN_T = Inches(0.9)
MARGIN_R = Inches(0.5)
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R


# ---------------------------------------------------------------------------
# python-pptx helpers
# ---------------------------------------------------------------------------

def _rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _fill_solid(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _no_fill(shape) -> None:
    shape.fill.background()


def _no_line(shape) -> None:
    shape.line.fill.background()


def _box(slide, left, top, width, height, fill=None, line_color=None, line_width_pt=0):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE
    if fill:
        _fill_solid(shape, fill)
    else:
        _no_fill(shape)
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width_pt or 1)
    else:
        _no_line(shape)
    return shape


def _text_in_rect(slide, rect: Rect, slide_w: int, slide_h: int, text: str, **kwargs) -> Any:
    return _text_box(
        slide, rect.left(slide_w), rect.top(slide_h),
        rect.width(slide_w), rect.height(slide_h), text, **kwargs
    )


def _text_box(slide, left, top, width, height, text: str, font_name="Red Hat Display",
              font_size=18, bold=False, italic=False, color=BLACK,
              align=PP_ALIGN.LEFT, word_wrap=True) -> Any:
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = word_wrap
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def _add_para(tf, text: str, font_name="Red Hat Text", font_size=13, bold=False,
              italic=False, color=BLACK, align=PP_ALIGN.LEFT, space_before=6) -> None:
    from pptx.util import Pt as _Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = _Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = _Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def _slide_bg(slide, color: RGBColor) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _detect_dark_bg(
    bp_layout: dict | None,
    img_dir: Path | None,
    fallback: bool = False,
) -> bool:
    """Sample the blueprint background image to decide if the slide is dark.

    Checks the text zone (y 30-45%, x 6-50%) where titles are placed.
    Falls back to ``fallback`` if no image can be analysed.
    """
    if not bp_layout or not img_dir:
        return fallback
    bg = bp_layout.get("background") or {}
    if bg.get("type") != "image" or not bg.get("image_ref"):
        return fallback
    img_path = img_dir / bg["image_ref"]
    if not img_path.exists():
        return fallback
    try:
        from PIL import Image
        import numpy as np
        img = Image.open(img_path).convert("RGB")
        arr = np.array(img)
        h, w = arr.shape[:2]
        region = arr[int(h * 0.30):int(h * 0.45), int(w * 0.06):int(w * 0.50)]
        brightness = float(region.mean())
        return brightness < 128
    except Exception:
        return fallback


def _text_color_for_bg(
    brand: "TemplateBrand",
    ctx: "SlideRenderContext | None",
    default_light: bool = False,
) -> RGBColor:
    """Pick white or black text based on whether the slide background is dark.

    In skeleton mode ``ctx.dark_bg`` is set from the variant (cover/section → dark).
    In blueprint/blank mode falls back to ``default_light`` (True → white text).
    """
    if ctx and ctx.dark_bg:
        return brand.white
    if ctx and not ctx.dark_bg:
        return brand.black
    return brand.white if default_light else brand.black


def _logo(slide, assets_root: Path, brand: "TemplateBrand | None" = None) -> None:
    """Skip logo placement when blueprint already baked it in."""
    if brand and brand.has_blueprint_bg:
        return  # logo is already part of the background image
    logo_path = assets_root / "logos" / "47c91d6a9ffd.png"
    if not logo_path.exists():
        candidates = list((assets_root / "logos").glob("*.png"))
        if not candidates:
            return
        logo_path = min(candidates, key=lambda p: p.stat().st_size)
    w = Inches(1.4)
    h = Inches(0.33)
    left = SLIDE_W - w - Inches(0.3)
    top  = Inches(0.18)
    try:
        slide.shapes.add_picture(str(logo_path), left, top, w, h)
    except Exception:
        pass


def _load_icon_migration(assets_root: Path) -> dict[str, str]:
    """Load old SHA-hash → new descriptive name mapping from icon_catalog.yaml."""
    catalog_path = assets_root / "icon_catalog.yaml"
    if not catalog_path.exists():
        return {}
    try:
        data = yaml.safe_load(catalog_path.read_text(encoding="utf-8")) or {}
        return data.get("id_migration", {})
    except Exception:
        return {}


_icon_migration_cache: dict[str, str] | None = None


def _place_icon(slide, icon_id: str, x, y, size, assets_root: Path) -> bool:
    """Place an icon PNG on the slide. Returns True if placed, False if not found.

    Supports both descriptive names (shield-security) and legacy SHA hashes,
    resolving old hashes via the id_migration map in icon_catalog.yaml.
    """
    global _icon_migration_cache
    icon_path = assets_root / "icons" / f"{icon_id}.png"
    if not icon_path.exists():
        if _icon_migration_cache is None:
            _icon_migration_cache = _load_icon_migration(assets_root)
        new_id = _icon_migration_cache.get(icon_id)
        if new_id:
            icon_path = assets_root / "icons" / f"{new_id}.png"
        if not icon_path.exists():
            return False
    try:
        slide.shapes.add_picture(str(icon_path), x, y, size, size)
        return True
    except Exception:
        return False


def _auto_font_size(text: str, base_pt: int, box_width_emu: int, box_height_emu: int) -> int:
    """Reduce font size when title text is long or multi-line."""
    n = len(text)
    lines = max(1, text.count("\n") + 1)
    if n <= 40:
        pt = base_pt
    elif n <= 70:
        pt = max(base_pt - 6, 18)
    elif n <= 100:
        pt = max(base_pt - 10, 16)
    else:
        pt = max(base_pt - 12, 14)
    # Fit within placeholder height (~1.15× line spacing)
    box_h_in = box_height_emu / 914400
    max_by_height = int(box_h_in / (lines * 1.15) * 72)
    if max_by_height > 0:
        pt = min(pt, max_by_height)
    return max(pt, 14)


def _estimate_line_count(text: str, box_width_emu: int, font_size_pt: int) -> int:
    """Estimate wrapped line count for stacked title-block layout."""
    char_width = max(int(Pt(font_size_pt) * 0.6), 1)
    chars_per_line = max(1, box_width_emu // char_width)
    num_lines = 0
    for segment in text.split("\n"):
        segment = segment or ""
        num_lines += max(1, (len(segment) + chars_per_line - 1) // chars_per_line)
    return max(1, num_lines)


_image_catalog_cache: dict[str, Any] | None = None


def _load_image_catalog(assets_root: Path) -> dict[str, Any]:
    global _image_catalog_cache
    if _image_catalog_cache is not None:
        return _image_catalog_cache
    catalog_path = assets_root / "image_catalog.yaml"
    if not catalog_path.exists():
        _image_catalog_cache = {}
        return _image_catalog_cache
    try:
        _image_catalog_cache = yaml.safe_load(catalog_path.read_text(encoding="utf-8")) or {}
    except Exception:
        _image_catalog_cache = {}
    return _image_catalog_cache


def _resolve_image(image_ref: str, assets_root: Path) -> Path | None:
    """Look up image_ref in image_catalog.yaml and return the file path."""
    if not image_ref:
        return None
    catalog = _load_image_catalog(assets_root)
    for entry in catalog.get("images", []):
        if entry.get("id") == image_ref:
            rel = entry.get("file", "")
            if rel:
                path = assets_root / rel
                if path.exists():
                    return path
    direct = assets_root / "images" / f"{image_ref}.png"
    if direct.exists():
        return direct
    for ext in ("jpg", "jpeg", "webp"):
        alt = assets_root / "images" / f"{image_ref}.{ext}"
        if alt.exists():
            return alt
    return None


def _layout_metrics(ctx: SlideRenderContext | None) -> LayoutMetrics:
    if ctx:
        return ctx.metrics(int(MARGIN_L), int(Inches(1.75)), int(CONTENT_W))
    return LayoutMetrics(
        margin_l=int(MARGIN_L),
        title_top=int(Inches(0.85)),
        title_height=int(Inches(0.75)),
        content_top=int(Inches(1.75)),
        content_width=int(CONTENT_W),
        content_height=int(Inches(5.0)),
        section_marker_top=int(Inches(0.12)),
    )


def _render_content_header(
    slide,
    title: str,
    brand: TemplateBrand,
    ctx: SlideRenderContext | None,
    ty: TypographyScale,
    section_marker: str | None = None,
) -> LayoutMetrics:
    """Stack section marker (optional), then title, then return content band metrics."""
    m = _layout_metrics(ctx)
    if section_marker:
        _text_box(
            slide, m.margin_l, m.section_marker_top, Inches(4), Inches(0.35),
            section_marker.upper(),
            font_name=brand.font_body, font_size=ty.section_marker,
            bold=True,
            color=brand.accent if brand.has_blueprint_bg else brand.mid_grey,
        )
    title_font_size = _auto_font_size(
        title, ty.slide_title, m.content_width, m.title_height,
    )
    _text_box(
        slide, m.margin_l, m.title_top, m.content_width, m.title_height,
        title, font_name=brand.font_heading, font_size=title_font_size,
        bold=True, color=brand.black,
    )
    return m


def _source_line(
    slide, text="Source: [cite]",
    brand: "TemplateBrand | None" = None,
    ctx: SlideRenderContext | None = None,
) -> None:
    b = brand or _DEFAULT_BRAND
    ty = ctx.typography if ctx else TypographyScale()
    m = _layout_metrics(ctx)
    source_top = m.content_top + m.content_height + int(Inches(0.05))
    source_top = min(source_top, int(SLIDE_H - Inches(0.4)))
    _text_box(
        slide, m.margin_l, source_top, m.content_width, Inches(0.3),
        text, font_name=b.font_body, font_size=ty.source, color=b.mid_grey,
    )


def _find_layout(master, name_hint: str | None = None):
    """Pick a slide layout by name hint, falling back to Blank then first layout."""
    if name_hint:
        hint = name_hint.lower()
        for layout in master.slide_layouts:
            if hint in layout.name.lower():
                return layout
    for layout in master.slide_layouts:
        if "blank" in layout.name.lower():
            return layout
    return master.slide_layouts[0]


def _add_empty_slide(
    prs: Presentation,
    master_idx: int = 0,
    layout_name_hint: str | None = "blank",
) -> Any:
    """Add a slide using Blank (or name hint) layout — avoids Title Slide placeholders."""
    try:
        master = prs.slide_masters[master_idx]
        layout = _find_layout(master, layout_name_hint)
        return prs.slides.add_slide(layout)
    except (IndexError, Exception):
        return prs.slides.add_slide(prs.slide_layouts[0])


def _read_placeholder_zones(slide, slide_w: int, slide_h: int) -> dict:
    """Read placeholder positions from a slide BEFORE clearing them.

    Returns a dict of placeholder type → (x_pct, y_pct, w_pct, h_pct).
    """
    zones: dict[str, dict] = {}
    for shape in list(slide.shapes):
        if not shape.is_placeholder:
            continue
        ph = shape.placeholder_format
        ptype = ph.type.name if ph.type else ""
        x_pct = round(shape.left / slide_w * 100, 2)
        y_pct = round(shape.top / slide_h * 100, 2)
        w_pct = round(shape.width / slide_w * 100, 2)
        h_pct = round(shape.height / slide_h * 100, 2)
        entry = {"type": ptype, "idx": ph.idx,
                 "x_pct": x_pct, "y_pct": y_pct, "w_pct": w_pct, "h_pct": h_pct}
        if ptype == "TITLE" and "TITLE" not in zones:
            zones["TITLE"] = entry
        elif ptype == "SUBTITLE":
            if "SUBTITLE" not in zones:
                zones["SUBTITLE"] = entry
        elif ptype in ("BODY", "OBJECT"):
            zones["BODY"] = entry
    return zones


def _ctx_from_placeholders(
    zones: dict,
    slide_w: int,
    slide_h: int,
    template_id: str | None,
    tmpl_yaml: Path | None,
) -> SlideRenderContext:
    """Build a SlideRenderContext from live placeholder measurements."""
    from generator.layout_context import (
        LayoutZones, Rect, TypographyScale, typography_for_template,
    )

    title_r = subtitle_r = body_r = presenter_r = None
    if "TITLE" in zones:
        z = zones["TITLE"]
        title_r = Rect(z["x_pct"], z["y_pct"], z["w_pct"], z["h_pct"])
    if "SUBTITLE" in zones:
        z = zones["SUBTITLE"]
        subtitle_r = Rect(z["x_pct"], z["y_pct"], z["w_pct"], z["h_pct"])
    if "BODY" in zones:
        z = zones["BODY"]
        body_r = Rect(z["x_pct"], z["y_pct"], z["w_pct"], z["h_pct"])

    # For layouts without a TITLE (e.g. CUSTOM_2 divider), synthesise a
    # title zone from the first SUBTITLE that isn't a right sidebar
    if title_r is None and subtitle_r is not None:
        if subtitle_r.x_pct < 50:
            title_r = Rect(
                subtitle_r.x_pct,
                max(subtitle_r.y_pct + subtitle_r.h_pct + 5.0, 25.0),
                min(subtitle_r.w_pct + 20.0, 65.0),
                20.0,
            )
            presenter_r = subtitle_r

    # Determine max text boundary — avoid right-side sidebar
    max_right = 92.0
    for z in zones.values():
        if z["x_pct"] > 70 and z["w_pct"] < 30:
            max_right = min(max_right, z["x_pct"] - 1.0)

    lz = LayoutZones(
        title=title_r,
        subtitle=subtitle_r,
        body=body_r,
        presenter=presenter_r,
        max_text_right_pct=max_right,
        max_text_bottom_pct=88.0,
    )
    ty = typography_for_template(template_id, tmpl_yaml)
    return SlideRenderContext(
        slide_w=slide_w, slide_h=slide_h,
        zones=lz, typography=ty,
    )


def _clear_placeholders(slide) -> None:
    """Remove all placeholder shapes inherited from the layout.

    When using a skeleton, layouts carry Title/Subtitle/Content placeholders
    that appear as empty boxes.  Deleting their XML elements leaves only the
    background artwork intact.
    """
    from pptx.oxml.ns import qn
    spTree = slide._element.find(qn("p:cSld")).find(qn("p:spTree"))
    if spTree is None:
        return
    to_remove = []
    for sp in spTree.findall(qn("p:sp")):
        nvSpPr = sp.find(qn("p:nvSpPr"))
        if nvSpPr is not None:
            nvPr = nvSpPr.find(qn("p:nvPr"))
            if nvPr is not None and nvPr.find(qn("p:ph")) is not None:
                to_remove.append(sp)
    for sp in to_remove:
        spTree.remove(sp)


# ---------------------------------------------------------------------------
# Skeleton + theme-manifest helpers
# ---------------------------------------------------------------------------

def _load_theme_manifest(skills_root: Path, template_id: str) -> dict | None:
    p = skills_root / "templates" / template_id / "theme-manifest.yaml"
    if not p.exists():
        return None
    with p.open() as f:
        return yaml.safe_load(f) or {}


def _layout_visual_color(layout) -> str | None:
    """Sample the background image of a layout to determine its dominant visual.

    Returns "RED", "DARK", "WHITE", or None (no image bg).
    Solid-fill colors in the XML are unreliable because background images
    override them visually.
    """
    NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

    blip = layout._element.find(f".//{{{NS_A}}}blip")
    if blip is None:
        # Fall back to solid fill
        import re as _re
        try:
            xml = layout._element.xml
            m = _re.search(
                r"<p:bg\b[^>]*>.*?<a:srgbClr val=\"([0-9A-Fa-f]{6})\"",
                xml, _re.DOTALL,
            )
            if m:
                h = m.group(1).upper()
                r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
                if r > 180 and g < 80:
                    return "RED"
                if r < 60 and g < 60 and b < 60:
                    return "DARK"
                if r > 200 and g > 200 and b > 200:
                    return "WHITE"
        except Exception:
            pass
        return None

    rid = blip.get(f"{{{NS_R}}}embed")
    if not rid or rid not in layout.part.rels:
        return None

    try:
        from PIL import Image
        import io
        blob = layout.part.rels[rid].target_part.blob
        img = Image.open(io.BytesIO(blob)).convert("RGB")
        w, h = img.size
        samples = [
            img.getpixel((w // 4, h // 4)),
            img.getpixel((w // 2, h // 2)),
            img.getpixel((3 * w // 4, h // 2)),
            img.getpixel((w // 2, 3 * h // 4)),
        ]
        avg_r = sum(p[0] for p in samples) // len(samples)
        avg_g = sum(p[1] for p in samples) // len(samples)
        avg_b = sum(p[2] for p in samples) // len(samples)

        if avg_r > 180 and avg_g < 80:
            return "RED"
        if avg_r < 60 and avg_g < 60 and avg_b < 60:
            return "DARK"
        if avg_r > 180 and avg_g > 180 and avg_b > 180:
            return "WHITE"
        return None
    except Exception:
        return None


# Slide variant → visual color expected from the background image
_VARIANT_VISUAL: dict[str, str] = {
    "cover":   "RED",
    "section": "DARK",
    "content": "WHITE",
}


def _build_layout_map(prs: Presentation, master_idx: int, variants: dict) -> dict:
    """
    Return {variant_name: layout} by sampling layout background images.

    Matches each variant to the first layout whose actual visual appearance
    (sampled from its background image) corresponds to the expected color.
    """
    result: dict = {}
    try:
        master = prs.slide_masters[master_idx]
    except IndexError:
        master = prs.slide_masters[0]

    for layout in master.slide_layouts:
        visual = _layout_visual_color(layout)
        if visual is None:
            continue
        for variant_name in variants:
            if variant_name in result:
                continue
            expected = _VARIANT_VISUAL.get(variant_name)
            if expected and visual == expected:
                result[variant_name] = layout
                break
        if len(result) >= len(variants):
            break

    fallback = master.slide_layouts[0]
    for name in variants:
        result.setdefault(name, fallback)
    return result


# Element → slide variant (determines which layout bg to inherit from skeleton)
_ELEMENT_VARIANT: dict[str, str] = {
    "title-block":         "cover",
    "divider":             "section",
    "closing":             "section",
    "agenda":              "content",
    "metric-card":         "content",
    "challenge-list":      "content",
    "tech-tile":           "content",
    "quote-block":         "content",
    "image-content":       "content",
    "data-table":          "content",
    "recommendation-card": "content",
    "timeline":            "content",
    "bar-chart":           "content",
}


# ---------------------------------------------------------------------------
# Blueprint replay — reproduce decorative shapes from extracted YAML
# ---------------------------------------------------------------------------

def _is_blueprint_placeholder(shape: dict) -> bool:
    """True when a blueprint shape is a text placeholder, not decoration."""
    kind = (shape.get("kind") or "").lower()
    if kind in ("placeholder", "ph"):
        return True
    if shape.get("is_placeholder") or shape.get("placeholder_type"):
        return True
    if shape.get("shape_type") in (14, "PLACEHOLDER", "MSO_SHAPE_TYPE.PLACEHOLDER"):
        return True
    return False


def _replay_blueprints(
    slide,
    layout_name: str,
    blueprints: dict,
    img_dir: Path,
    *,
    skip_text_shapes: bool = False,
) -> None:
    """
    Place every decorative shape recorded in the blueprints YAML onto the slide.

    When the layout has no background of its own, we first replay the parent
    slide master's decorative shapes (red stripe, logo, etc.) so that the
    master branding is always present.

    Placeholder shapes and optional text boxes (e.g. closing slide content)
    are skipped — renderers draw that text themselves.
    """
    layouts = blueprints.get("layouts", {})
    layout  = layouts.get(layout_name)
    if not layout:
        return

    W = int(prs_width_emu_cache[0]) if prs_width_emu_cache else 12192000
    H = int(prs_height_emu_cache[0]) if prs_height_emu_cache else 6858000

    bg = layout.get("background")
    master_idx = layout.get("master", 0)
    master_key = f"_MASTER_{master_idx}"
    master = layouts.get(master_key)

    # --- Step 1a: If layout has no bg, replay master bg + master shapes first ---
    if not bg or bg.get("type") not in ("image", "solid"):
        if master:
            m_bg = master.get("background")
            if m_bg and m_bg.get("type") == "image":
                bg_ref = m_bg.get("image_ref")
                if bg_ref:
                    bg_path = img_dir / bg_ref
                    if bg_path.exists():
                        try:
                            slide.shapes.add_picture(str(bg_path), 0, 0, W, H)
                        except Exception:
                            pass
            elif m_bg and m_bg.get("type") == "solid":
                try:
                    _slide_bg(slide, _rgb(m_bg["color"].lstrip("#")))
                except Exception:
                    pass
            # Replay master decorative shapes
            for shape in master.get("decorative_shapes", []):
                if _is_blueprint_placeholder(shape):
                    continue
                _replay_single_shape(
                    slide, shape, W, H, img_dir, skip_text_shapes=skip_text_shapes,
                )
    else:
        # --- Step 1b: Layout has its own background ---
        if bg.get("type") == "image":
            bg_ref = bg.get("image_ref")
            if bg_ref:
                bg_path = img_dir / bg_ref
                if bg_path.exists():
                    try:
                        slide.shapes.add_picture(str(bg_path), 0, 0, W, H)
                    except Exception:
                        pass
        elif bg.get("type") == "solid":
            try:
                _slide_bg(slide, _rgb(bg["color"].lstrip("#")))
            except Exception:
                pass

    # --- Step 2: Overlay layout's own decorative shapes ---
    for shape in layout.get("decorative_shapes", []):
        if _is_blueprint_placeholder(shape):
            continue
        _replay_single_shape(
            slide, shape, W, H, img_dir, skip_text_shapes=skip_text_shapes,
        )


def _replay_single_shape(
    slide,
    shape: dict,
    W: int,
    H: int,
    img_dir: Path,
    *,
    skip_text_shapes: bool = False,
) -> None:
    """Draw a single decorative shape from a blueprint dict onto a slide."""
    kind = shape.get("kind")
    if _is_blueprint_placeholder(shape):
        return
    if skip_text_shapes and kind == "textbox":
        return
    x = int(shape.get("x_pct", 0) / 100 * W)
    y = int(shape.get("y_pct", 0) / 100 * H)
    w = int(shape.get("w_pct", 0) / 100 * W)
    h = int(shape.get("h_pct", 0) / 100 * H)

    try:
        if kind == "picture":
            img_ref = shape.get("image_ref")
            if not img_ref:
                return
            img_path = img_dir / img_ref
            if img_path.exists():
                slide.shapes.add_picture(str(img_path), x, y, w, h)

        elif kind == "line":
            li = shape.get("line", {})
            color_hex = li.get("color", "none")
            width_pt  = li.get("width_pt", 1)
            from pptx.util import Pt as _Pt
            from pptx.enum.shapes import MSO_CONNECTOR_TYPE
            connector = slide.shapes.add_connector(
                MSO_CONNECTOR_TYPE.STRAIGHT,
                x, y, x + w, y + h
            )
            if color_hex and color_hex != "none":
                connector.line.color.rgb = _rgb(color_hex.lstrip("#") if not color_hex.startswith("#") else color_hex[1:])
            connector.line.width = _Pt(width_pt)
            dash = li.get("dash", "solid")
            if dash and dash != "solid":
                try:
                    from pptx.oxml.ns import qn
                    from lxml import etree
                    ln_el = connector.line._ln
                    prstDash = etree.SubElement(ln_el, qn("a:prstDash"))
                    dash_map_rev = {
                        "dash": "dash", "dash-dot": "dashDot",
                        "dot": "dot", "long-dash": "lgDash",
                        "long-dash-dot": "lgDashDot",
                    }
                    prstDash.set("val", dash_map_rev.get(dash, "dash"))
                except Exception:
                    pass

        elif kind == "shape":
            fill = shape.get("fill", {})
            li   = shape.get("line", {})
            rect = _box(slide, x, y, w, h)
            if fill.get("type") == "solid" and fill.get("color"):
                _fill_solid(rect, _rgb(fill["color"].lstrip("#")))
            else:
                _no_fill(rect)
            line_color = li.get("color", "none")
            if line_color and line_color != "none":
                rect.line.color.rgb = _rgb(line_color.lstrip("#"))
                rect.line.width = Pt(li.get("width_pt", 1))
            else:
                _no_line(rect)

        elif kind == "textbox":
            txt   = shape.get("text_sample", "")
            font  = shape.get("font_name") or "Red Hat Text"
            sz    = shape.get("font_size_pt") or 11
            color = shape.get("font_color", "#000000")
            try:
                clr = _rgb(color.lstrip("#"))
            except Exception:
                clr = BLACK
            if txt:
                _text_box(slide, x, y, w, h, txt, font_name=font,
                          font_size=sz, color=clr)

    except Exception:
        pass


# Module-level slide dimension cache (set by build_deck before rendering)
prs_width_emu_cache:  list[int] = []
prs_height_emu_cache: list[int] = []


def _place_hero_image(
    slide, img_path: Path, slide_w: int, slide_h: int,
    max_w_pct: float = 38.0, max_h_pct: float = 70.0,
    x_pct: float = 58.0, y_pct: float = 15.0,
) -> None:
    """Place an image on the right side of a slide preserving aspect ratio."""
    try:
        from PIL import Image as PILImage
        pil = PILImage.open(img_path)
        iw, ih = pil.size
    except Exception:
        return
    if iw < 200 or ih < 200:
        return
    box_w = int(max_w_pct / 100 * slide_w)
    box_h = int(max_h_pct / 100 * slide_h)
    scale = min(box_w / iw, box_h / ih)
    final_w = int(iw * scale)
    final_h = int(ih * scale)
    left = int(x_pct / 100 * slide_w) + (box_w - final_w) // 2
    top = int(y_pct / 100 * slide_h) + (box_h - final_h) // 2
    try:
        slide.shapes.add_picture(str(img_path), left, top, final_w, final_h)
    except Exception:
        pass


# ---------------------------------------------------------------------------
# Element renderers  (all accept brand: TemplateBrand)
# ---------------------------------------------------------------------------

def render_title_block(
    slide, spec: dict, assets_root: Path,
    brand: TemplateBrand = None, ctx: SlideRenderContext | None = None,
) -> None:
    b = brand or _DEFAULT_BRAND
    ty = ctx.typography if ctx else TypographyScale()
    if not b.has_blueprint_bg:
        _slide_bg(slide, b.dark)
    _logo(slide, assets_root, b)
    title_color = _text_color_for_bg(b, ctx, default_light=True)

    if ctx and ctx.zones.title:
        title_text = spec.get("title", "Presentation Title")
        title_rect = ctx.zones.title
        max_text_right = min(ctx.zones.max_text_right_pct, 55.0)
        max_w_pct = min(title_rect.w_pct, max_text_right - title_rect.x_pct)
        # Title slide: don't let placeholder zones push headline below ~35% of slide height
        title_y_pct = min(title_rect.y_pct, 35.0)
        clamped_title = Rect(title_rect.x_pct, title_y_pct, max_w_pct, title_rect.h_pct)

        title_font_size = _auto_font_size(
            title_text, ty.deck_title,
            clamped_title.width(ctx.slide_w), clamped_title.height(ctx.slide_h),
        )
        _text_in_rect(
            slide, clamped_title, ctx.slide_w, ctx.slide_h,
            title_text,
            font_name=b.font_heading, font_size=title_font_size, bold=True, color=title_color,
        )

        num_lines = _estimate_line_count(
            title_text, clamped_title.width(ctx.slide_w), title_font_size,
        )
        title_used_h_pct = num_lines * title_font_size * 1.4 / (ctx.slide_h / 914400 * 72) * 100

        sub_y = max(
            clamped_title.y_pct + title_used_h_pct + 1.0,
            (ctx.zones.subtitle.y_pct if ctx.zones.subtitle else clamped_title.y_pct + clamped_title.h_pct + 1.0),
        )
        sub_y = min(sub_y, 75.0)

        if spec.get("subheading"):
            sub_w = min(
                ctx.zones.subtitle.w_pct if ctx.zones.subtitle else max_w_pct,
                max_w_pct,
            )
            sub_h = ctx.zones.subtitle.h_pct if ctx.zones.subtitle else 6.0
            sub_rect = Rect(clamped_title.x_pct, sub_y, sub_w, sub_h)
            sub_font = _auto_font_size(
                spec["subheading"], ty.deck_subtitle,
                sub_rect.width(ctx.slide_w), sub_rect.height(ctx.slide_h),
            )
            _text_in_rect(
                slide, sub_rect, ctx.slide_w, ctx.slide_h, spec["subheading"],
                font_name=b.font_body, font_size=sub_font, color=title_color,
            )
            sub_lines = _estimate_line_count(
                spec["subheading"], sub_rect.width(ctx.slide_w), sub_font,
            )
            sub_y += sub_lines * sub_font * 1.3 / (ctx.slide_h / 914400 * 72) * 100 + 0.5

        bottom = spec.get("presenter", "")
        if spec.get("date"):
            bottom += f"  |  {spec['date']}" if bottom else spec["date"]
        if bottom.strip():
            pres_rect = Rect(clamped_title.x_pct, sub_y + 1.0, max_w_pct, 4.0)
            pres_font = _auto_font_size(
                bottom, ty.deck_presenter,
                pres_rect.width(ctx.slide_w), pres_rect.height(ctx.slide_h),
            )
            _text_in_rect(
                slide, pres_rect, ctx.slide_w, ctx.slide_h, bottom,
                font_name=b.font_body, font_size=pres_font, color=title_color,
            )

        pass  # title-block relies on blueprint background as brand visual
    else:
        _text_box(slide, MARGIN_L, Inches(2.2), SLIDE_W * 0.6, Inches(1.8),
                  spec.get("title", "Presentation Title"),
                  font_name=b.font_heading, font_size=ty.deck_title, bold=True, color=title_color)
        if spec.get("subheading"):
            _text_box(slide, MARGIN_L, Inches(4.1), SLIDE_W * 0.6, Inches(0.6),
                      spec["subheading"], font_name=b.font_body, font_size=ty.deck_subtitle,
                      color=title_color)


def render_divider(
    slide, spec: dict, assets_root: Path,
    brand: TemplateBrand = None, ctx: SlideRenderContext | None = None,
) -> None:
    b = brand or _DEFAULT_BRAND
    ty = ctx.typography if ctx else TypographyScale()
    if not b.has_blueprint_bg:
        _slide_bg(slide, b.primary)
    _logo(slide, assets_root, b)
    txt_color = _text_color_for_bg(b, ctx, default_light=True)

    if ctx and ctx.zones.title:
        if spec.get("section_marker"):
            marker = Rect(
                ctx.zones.title.x_pct,
                max(2.0, ctx.zones.title.y_pct - 8.0),
                min(45.0, ctx.zones.title.w_pct),
                6.0,
            )
            _text_in_rect(
                slide, marker, ctx.slide_w, ctx.slide_h,
                spec["section_marker"].upper(),
                font_name=b.font_body, font_size=ty.divider_marker,
                bold=True, color=b.accent if b.has_blueprint_bg else txt_color,
            )
        headline = spec.get("headline", "Section")
        title_rect = ctx.zones.title
        headline_font_size = _auto_font_size(
            headline, ty.divider_headline,
            title_rect.width(ctx.slide_w), title_rect.height(ctx.slide_h),
        )
        _text_in_rect(
            slide, title_rect, ctx.slide_w, ctx.slide_h,
            headline,
            font_name=b.font_heading, font_size=headline_font_size, bold=True, color=txt_color,
        )
    else:
        if spec.get("section_marker"):
            _text_box(slide, MARGIN_L, Inches(2.4), SLIDE_W * 0.7, Inches(0.4),
                      spec["section_marker"].upper(), font_name=b.font_body,
                      font_size=ty.divider_marker, bold=True, color=txt_color)
        _text_box(slide, MARGIN_L, Inches(2.85), SLIDE_W * 0.7, Inches(2.0),
                  spec.get("headline", "Section"), font_name=b.font_heading,
                  font_size=ty.divider_headline, bold=True, color=txt_color)


def render_agenda(
    slide, spec: dict, assets_root: Path,
    brand: TemplateBrand = None, ctx: SlideRenderContext | None = None,
) -> None:
    b = brand or _DEFAULT_BRAND
    ty = ctx.typography if ctx else TypographyScale()
    if not b.has_blueprint_bg:
        _slide_bg(slide, b.white)
    _logo(slide, assets_root, b)
    m = _render_content_header(slide, spec.get("title", "Agenda"), b, ctx, ty)
    items = spec.get("items", [])[:6]
    n = max(len(items), 1)
    row_h = min(int(m.content_height / n * 0.88), int(Inches(0.95)))
    gap = int(Inches(0.08))
    for i, item in enumerate(items):
        y = m.content_top + i * (row_h + gap)
        _box(slide, m.margin_l, y, Pt(3), row_h, fill=b.primary)
        _text_box(slide, m.margin_l + Pt(10), y + int(Inches(0.1)), Inches(0.45), Inches(0.4),
                  str(i + 1), font_name=b.font_heading, font_size=ty.agenda_topic,
                  bold=True, color=b.primary, align=PP_ALIGN.CENTER)
        topic_left = m.margin_l + Inches(0.65)
        text_w = m.content_width - Inches(0.7)
        icon_id = item.get("icon")
        if icon_id:
            icon_size = Inches(0.35)
            icon_y = y + (row_h - icon_size) // 2
            if _place_icon(slide, icon_id, m.margin_l + Inches(0.52), icon_y,
                           icon_size, assets_root):
                topic_left = m.margin_l + Inches(0.52) + icon_size + Inches(0.08)
                text_w = m.content_width - (topic_left - m.margin_l) - Inches(0.05)
        _text_box(slide, topic_left, y + int(Inches(0.04)),
                  text_w, int(row_h * 0.42),
                  item.get("topic", ""), font_name=b.font_heading,
                  font_size=ty.agenda_topic, bold=True, color=b.black)
        if item.get("detail"):
            _text_box(slide, topic_left, y + int(row_h * 0.46),
                      text_w, int(row_h * 0.48),
                      item["detail"], font_name=b.font_body,
                      font_size=ty.agenda_detail, color=b.dark_grey)


def _metric_text(value: Any, default: str = "") -> str:
    """Coerce a metric field to display text; LLM output may nest dicts."""
    if isinstance(value, dict):
        return str(value.get("value", value))
    if value is None:
        return default
    return str(value)


def render_metric_cards(
    slide, spec: dict, assets_root: Path,
    brand: TemplateBrand = None, ctx: SlideRenderContext | None = None,
) -> None:
    b = brand or _DEFAULT_BRAND
    ty = ctx.typography if ctx else TypographyScale()
    if not b.has_blueprint_bg:
        _slide_bg(slide, b.white)
    _logo(slide, assets_root, b)
    m = _render_content_header(
        slide, spec.get("title", "Key Metrics"), b, ctx, ty,
        section_marker=spec.get("section_marker"),
    )
    cards = (spec.get("cards") or spec.get("metrics") or [])[:4]
    n = max(len(cards), 1)
    gap = int(Inches(0.25))
    card_w = (m.content_width - gap * (n - 1)) // n
    card_h = min(int(m.content_height * 0.85), int(Inches(2.6)))
    for i, card in enumerate(cards):
        x = m.margin_l + i * (card_w + gap)
        y = m.content_top
        _box(slide, x, y, card_w, card_h, fill=b.white, line_color=b.light_bg, line_width_pt=1)
        _box(slide, x, y, Pt(4), card_h, fill=b.primary)
        if isinstance(card, dict):
            val_text = _metric_text(card.get("value", ""))
            label_text = _metric_text(card.get("label", ""))
        else:
            val_text = _metric_text(card)
            label_text = ""
        val_is_stat = (
            len(val_text.replace("\n", "")) <= 8
            and any(c.isdigit() or c in "%€$£x×" for c in val_text)
        )
        val_font_size = ty.metric_value if val_is_stat else ty.card_headline
        val_bold = True
        val_color = b.primary if val_is_stat else b.black
        _text_box(slide, x + Inches(0.15), y + Inches(0.25),
                  card_w - Inches(0.2), int(card_h * 0.38),
                  val_text, font_name=b.font_heading,
                  font_size=val_font_size, bold=val_bold, color=val_color)
        _text_box(slide, x + Inches(0.15), y + int(card_h * 0.42),
                  card_w - Inches(0.2), int(card_h * 0.52),
                  label_text, font_name=b.font_body,
                  font_size=ty.metric_label, color=b.dark_grey, word_wrap=True)
        icon_id = card.get("icon")
        if icon_id:
            icon_size = Inches(0.3)
            _place_icon(slide, icon_id, x + card_w - icon_size - Inches(0.1),
                        y + Inches(0.1), icon_size, assets_root)
    _source_line(slide, spec.get("source", "Source: [cite]"), b, ctx)


def render_challenge_list(
    slide, spec: dict, assets_root: Path,
    brand: TemplateBrand = None, ctx: SlideRenderContext | None = None,
) -> None:
    b = brand or _DEFAULT_BRAND
    ty = ctx.typography if ctx else TypographyScale()
    if not b.has_blueprint_bg:
        _slide_bg(slide, b.white)
    _logo(slide, assets_root, b)
    m = _render_content_header(
        slide, spec.get("title", "Customer Challenges"), b, ctx, ty,
        section_marker=spec.get("section_marker"),
    )
    challenges = spec.get("challenges", [])[:4]
    n = len(challenges)

    if n <= 3:
        row_h = min(int(m.content_height / max(n, 1) * 0.88), int(Inches(1.35)))
        gap = int(Inches(0.1))
        for i, ch in enumerate(challenges):
            y = m.content_top + i * (row_h + gap)
            badge = Inches(0.5)
            icon_id = ch.get("icon")
            icon_size = Inches(0.4)
            if not (icon_id and _place_icon(
                slide, icon_id, m.margin_l, y + (badge - icon_size) // 2,
                icon_size, assets_root,
            )):
                _box(slide, m.margin_l, y, badge, badge, fill=b.primary)
                _text_box(slide, m.margin_l, y, badge, badge, str(i + 1),
                          font_name=b.font_heading, font_size=ty.card_headline,
                          bold=True, color=b.white, align=PP_ALIGN.CENTER)
            _text_box(slide, m.margin_l + Inches(0.65), y,
                      m.content_width - Inches(0.7), int(row_h * 0.4),
                      ch.get("headline", ""), font_name=b.font_heading,
                      font_size=ty.card_headline, bold=True, color=b.black)
            _text_box(slide, m.margin_l + Inches(0.65), y + int(row_h * 0.42),
                      m.content_width - Inches(0.7), int(row_h * 0.55),
                      ch.get("body", ""), font_name=b.font_body,
                      font_size=ty.card_body, color=b.dark_grey, word_wrap=True)
    else:
        gap = int(Inches(0.25))
        card_w = (m.content_width - gap) // 2
        card_h = min(int(m.content_height / 2 * 0.9), int(Inches(2.2)))
        for i, ch in enumerate(challenges):
            col, row = i % 2, i // 2
            x = m.margin_l + col * (card_w + gap)
            y = m.content_top + row * (card_h + gap)
            _box(slide, x, y, card_w, card_h, fill=b.light_bg)
            _box(slide, x, y, card_w, Pt(4), fill=b.primary)
            _text_box(slide, x + Inches(0.15), y + Inches(0.12),
                      card_w - Inches(0.3), Inches(0.45),
                      ch.get("headline", ""), font_name=b.font_heading,
                      font_size=ty.card_headline, bold=True, color=b.black)
            _text_box(slide, x + Inches(0.15), y + Inches(0.58),
                      card_w - Inches(0.3), card_h - Inches(0.65),
                      ch.get("body", ""), font_name=b.font_body,
                      font_size=ty.card_body, color=b.dark_grey, word_wrap=True)


def render_tech_tiles(
    slide, spec: dict, assets_root: Path,
    brand: TemplateBrand = None, ctx: SlideRenderContext | None = None,
) -> None:
    b = brand or _DEFAULT_BRAND
    ty = ctx.typography if ctx else TypographyScale()
    if not b.has_blueprint_bg:
        _slide_bg(slide, b.white)
    _logo(slide, assets_root, b)
    m = _render_content_header(
        slide, spec.get("title", "Solution Portfolio"), b, ctx, ty,
        section_marker=spec.get("section_marker"),
    )
    tiles = spec.get("tiles", [])[:6]
    cols = 3 if len(tiles) > 4 else 2
    rows = (len(tiles) + cols - 1) // cols
    gap_x, gap_y = int(Inches(0.15)), int(Inches(0.12))
    tile_w = (m.content_width - gap_x * (cols - 1)) // cols
    tile_h = (m.content_height - gap_y * (rows - 1)) // rows
    pillar_colors = {"Technology": b.primary, "Operations": b.accent, "Assurance": b.dark_grey}
    for i, tile in enumerate(tiles):
        col, row = i % cols, i // cols
        x = m.margin_l + col * (tile_w + gap_x)
        y = m.content_top + row * (tile_h + gap_y)
        strip = pillar_colors.get(tile.get("pillar", ""), b.primary)
        _box(slide, x, y, tile_w, tile_h, fill=b.white, line_color=b.light_bg, line_width_pt=1)
        _box(slide, x, y, tile_w, Pt(5), fill=strip)
        if tile.get("pillar"):
            _text_box(slide, x + Inches(0.1), y + Inches(0.06),
                      tile_w - Inches(0.15), Inches(0.22),
                      tile["pillar"].upper(), font_name=b.font_body,
                      font_size=8, bold=True, color=strip)
        name_top = y + Inches(0.32)
        icon_id = tile.get("icon")
        if icon_id:
            icon_size = Inches(0.5)
            icon_x = x + (tile_w - icon_size) // 2
            icon_y = y + Inches(0.28) if tile.get("pillar") else y + Inches(0.15)
            if _place_icon(slide, icon_id, icon_x, icon_y, icon_size, assets_root):
                name_top = icon_y + icon_size + Inches(0.06)
        _text_box(slide, x + Inches(0.1), name_top,
                  tile_w - Inches(0.15), int(tile_h * 0.3),
                  tile.get("name", ""), font_name=b.font_heading,
                  font_size=ty.card_headline, bold=True, color=b.black)
        desc_top = name_top + int(tile_h * 0.38) - int(y + Inches(0.32))
        _text_box(slide, x + Inches(0.1), desc_top,
                  tile_w - Inches(0.15), y + tile_h - desc_top - int(Inches(0.05)),
                  tile.get("description", ""), font_name=b.font_body,
                  font_size=ty.card_body, color=b.dark_grey, word_wrap=True)


def render_quote_block(
    slide, spec: dict, assets_root: Path,
    brand: TemplateBrand = None, ctx: SlideRenderContext | None = None,
) -> None:
    b = brand or _DEFAULT_BRAND
    ty = ctx.typography if ctx else TypographyScale()
    if not b.has_blueprint_bg:
        _slide_bg(slide, b.white)
    _logo(slide, assets_root, b)
    m = _render_content_header(
        slide, spec.get("title", ""), b, ctx, ty,
        section_marker=spec.get("section_marker"),
    )
    _text_box(slide, m.margin_l, m.content_top, int(m.content_width * 0.65), int(m.content_height * 0.72),
              spec.get("quote", ""), font_name=b.font_heading, font_size=ty.quote,
              italic=True, color=b.black, word_wrap=True)
    if spec.get("attribution"):
        _text_box(slide, m.margin_l, m.content_top + int(m.content_height * 0.74),
                  int(m.content_width * 0.65), int(m.content_height * 0.2),
                  spec["attribution"], font_name=b.font_body, font_size=ty.body,
                  bold=True, color=b.primary)


def render_image_content(
    slide, spec: dict, assets_root: Path,
    brand: TemplateBrand = None, ctx: SlideRenderContext | None = None,
) -> None:
    b = brand or _DEFAULT_BRAND
    ty = ctx.typography if ctx else TypographyScale()
    if not b.has_blueprint_bg:
        _slide_bg(slide, b.white)
    _logo(slide, assets_root, b)
    m = _render_content_header(
        slide, spec.get("title", ""), b, ctx, ty,
        section_marker=spec.get("section_marker"),
    )
    text_w = int(m.content_width * 0.52)
    img_w = m.content_width - text_w - int(Inches(0.2))
    img_x = m.margin_l + text_w + int(Inches(0.15))
    img_path = None
    if spec.get("image_ref"):
        img_path = _resolve_image(spec["image_ref"], assets_root)
    elif spec.get("image_path"):
        candidate = Path(spec["image_path"])
        img_path = candidate if candidate.exists() else None

    if img_path:
        try:
            from PIL import Image as PILImage
            pil = PILImage.open(img_path)
            iw, ih = pil.size
            scale = min(img_w / iw, m.content_height / ih)
            fw, fh = int(iw * scale), int(ih * scale)
            cx = img_x + (img_w - fw) // 2
            cy = m.content_top + (m.content_height - fh) // 2
            slide.shapes.add_picture(str(img_path), cx, cy, fw, fh)
        except Exception:
            img_path = None

    if not img_path:
        _box(slide, img_x, m.content_top, img_w, m.content_height, fill=b.light_bg)
        _text_box(slide, img_x + Inches(0.15), m.content_top + m.content_height // 2 - Inches(0.25),
                  img_w - Inches(0.3), Inches(0.5),
                  "[Replace with approved image]",
                  font_name=b.font_body, font_size=ty.body_small, color=b.mid_grey, align=PP_ALIGN.CENTER)
    bullets = spec.get("bullets", [])[:6]
    row_h = min(int(m.content_height / max(len(bullets), 1)), int(Inches(0.9)))
    for i, bullet in enumerate(bullets):
        y = m.content_top + i * row_h
        if bullet.get("headline"):
            _text_box(slide, m.margin_l + Inches(0.2), y, text_w - Inches(0.3), int(row_h * 0.38),
                      bullet["headline"], font_name=b.font_heading, font_size=ty.card_headline,
                      bold=True, color=b.black)
        if bullet.get("body"):
            _text_box(slide, m.margin_l + Inches(0.2), y + int(row_h * 0.4), text_w - Inches(0.3),
                      int(row_h * 0.55), bullet["body"], font_name=b.font_body,
                      font_size=ty.card_body, color=b.dark_grey, word_wrap=True)


def render_data_table(
    slide, spec: dict, assets_root: Path,
    brand: TemplateBrand = None, ctx: SlideRenderContext | None = None,
) -> None:
    from pptx.util import Pt as _Pt
    b = brand or _DEFAULT_BRAND
    ty = ctx.typography if ctx else TypographyScale()
    if not b.has_blueprint_bg:
        _slide_bg(slide, b.white)
    _logo(slide, assets_root, b)
    m = _render_content_header(
        slide, spec.get("title", "Comparison"), b, ctx, ty,
        section_marker=spec.get("section_marker"),
    )
    headers, rows, takeaway = spec.get("headers", []), spec.get("rows", []), spec.get("takeaway", "")
    if not headers:
        return

    table_top = m.content_top
    table_left = m.margin_l
    table_w = m.content_width
    table_h_budget = int(m.content_height * 0.72)
    row_h = min(Inches(0.55), table_h_budget // max(len(rows) + 1, 1))
    tbl = slide.shapes.add_table(
        len(rows) + 1, len(headers), table_left, table_top,
        table_w, row_h * (len(rows) + 1),
    ).table
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = b.primary
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.name = b.font_heading
                run.font.size = _Pt(ty.card_headline)
                run.font.bold = True
                run.font.color.rgb = b.white
    for ri, row in enumerate(rows):
        row_bg = b.light_bg if ri % 2 == 0 else b.white
        for ci, val in enumerate(row[:len(headers)]):
            cell = tbl.cell(ri + 1, ci)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_bg
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = b.font_body
                    run.font.size = _Pt(ty.card_body)
                    run.font.color.rgb = b.black
    if takeaway:
        callout_y = table_top + row_h * (len(rows) + 1) + Inches(0.1)
        callout_h = max(int(table_h_budget * 0.22), Inches(0.45))
        _box(slide, table_left, callout_y, table_w, callout_h,
             fill=None, line_color=b.primary, line_width_pt=1.5)
        _text_box(slide, table_left + Inches(0.15), callout_y + Inches(0.08),
                  table_w - Inches(0.3), max(callout_h - Inches(0.15), Inches(0.35)),
                  takeaway, font_name=b.font_body, font_size=ty.body_small,
                  italic=True, color=b.dark_grey, word_wrap=True)
    _source_line(slide, spec.get("source", "Source: [cite]"), b, ctx)


def render_recommendation_cards(
    slide, spec: dict, assets_root: Path,
    brand: TemplateBrand = None, ctx: SlideRenderContext | None = None,
) -> None:
    b = brand or _DEFAULT_BRAND
    ty = ctx.typography if ctx else TypographyScale()
    if not b.has_blueprint_bg:
        _slide_bg(slide, b.white)
    _logo(slide, assets_root, b)
    m = _render_content_header(slide, spec.get("title", "Where to Start"), b, ctx, ty)
    cards = spec.get("cards", [])[:4]
    n = max(len(cards), 1)
    gap = int(Inches(0.25))
    card_w = (m.content_width - gap * (n - 1)) // n
    card_h = min(int(m.content_height * 0.75), int(Inches(3.2)))
    for i, card in enumerate(cards):
        x = m.margin_l + i * (card_w + gap)
        y = m.content_top
        _box(slide, x, y, card_w, card_h, fill=b.white, line_color=b.light_bg, line_width_pt=1)
        badge = Inches(0.5)
        badge_x = x + Inches(0.15)
        badge_y = y + Inches(0.15)
        icon_id = card.get("icon")
        icon_size = Inches(0.35)
        if not (icon_id and _place_icon(
            slide, icon_id,
            badge_x + (badge - icon_size) // 2,
            badge_y + (badge - icon_size) // 2,
            icon_size, assets_root,
        )):
            _box(slide, badge_x, badge_y, badge, badge, fill=b.primary)
            _text_box(slide, badge_x, badge_y, badge, badge, str(i + 1),
                      font_name=b.font_heading, font_size=ty.agenda_topic,
                      bold=True, color=b.white, align=PP_ALIGN.CENTER)
        _text_box(slide, x + Inches(0.12), y + Inches(0.85), card_w - Inches(0.24), Inches(0.6),
                  card.get("headline", ""), font_name=b.font_heading,
                  font_size=ty.card_headline, bold=True, color=b.black, word_wrap=True)
        _text_box(slide, x + Inches(0.12), y + Inches(1.5), card_w - Inches(0.24), card_h - Inches(1.65),
                  card.get("body", ""), font_name=b.font_body,
                  font_size=ty.card_body, color=b.dark_grey, word_wrap=True)
    if spec.get("cta"):
        _text_box(slide, m.margin_l, m.content_top + card_h + int(Inches(0.12)),
                  m.content_width, Inches(0.4), spec["cta"],
                  font_name=b.font_body, font_size=ty.body, italic=True,
                  color=b.primary, align=PP_ALIGN.CENTER)


def render_closing(
    slide, spec: dict, assets_root: Path,
    brand: TemplateBrand = None, ctx: SlideRenderContext | None = None,
) -> None:
    b = brand or _DEFAULT_BRAND
    ty = ctx.typography if ctx else TypographyScale()
    boilerplate = spec.get("boilerplate")

    if b.has_blueprint_bg:
        if not boilerplate:
            return
        slide_w = ctx.slide_w if ctx else int(SLIDE_W)
        slide_h = ctx.slide_h if ctx else int(SLIDE_H)
        bp_color = _text_color_for_bg(b, ctx, default_light=False)
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            top_pct = shape.top / slide_h * 100 if shape.top else 0
            if 45.0 <= top_pct <= 55.0:
                sample = shape.text_frame.text.strip().lower()
                if sample and ("open source" in sample or "red hat is" in sample):
                    shape.text_frame.clear()
        bp_rect = Rect(1.4, 48.0, 38.0, 15.0)
        _text_in_rect(
            slide, bp_rect, slide_w, slide_h, boilerplate,
            font_name=b.font_body, font_size=11, color=bp_color, word_wrap=True,
        )
        return

    txt = _text_color_for_bg(b, ctx, default_light=True)

    if not b.has_blueprint_bg:
        _slide_bg(slide, b.dark)
    _logo(slide, assets_root, b)
    boilerplate = boilerplate or (
        "Red Hat is the world's leading provider of enterprise open source software solutions."
    )
    footer = "linkedin.com/company/red-hat"

    if ctx and ctx.zones.title:
        title_rect = ctx.zones.title
        _text_in_rect(
            slide, title_rect, ctx.slide_w, ctx.slide_h, "Thank you.",
            font_name=b.font_heading, font_size=ty.closing_title, bold=True, color=txt,
        )
        body_rect = Rect(
            title_rect.x_pct,
            title_rect.y_pct + title_rect.h_pct + 3.0,
            title_rect.w_pct,
            18.0,
        )
        _text_in_rect(
            slide, body_rect, ctx.slide_w, ctx.slide_h, boilerplate,
            font_name=b.font_body, font_size=ty.body, color=txt, word_wrap=True,
        )
        footer_rect = Rect(
            title_rect.x_pct,
            title_rect.y_pct + title_rect.h_pct + 24.0,
            title_rect.w_pct,
            8.0,
        )
        _text_in_rect(
            slide, footer_rect, ctx.slide_w, ctx.slide_h, footer,
            font_name=b.font_body, font_size=ty.body_small, color=txt,
        )
    else:
        _text_box(slide, MARGIN_L, Inches(2.0), SLIDE_W * 0.65, Inches(1.2),
                  "Thank you.", font_name=b.font_heading, font_size=ty.closing_title,
                  bold=True, color=txt)
        _text_box(slide, MARGIN_L, Inches(3.6), SLIDE_W * 0.65, Inches(1.5),
                  boilerplate, font_name=b.font_body, font_size=ty.body, color=txt,
                  word_wrap=True)
        _text_box(slide, MARGIN_L, Inches(5.4), SLIDE_W * 0.65, Inches(0.4),
                  footer, font_name=b.font_body, font_size=ty.body_small, color=txt)


def render_timeline(
    slide, spec: dict, assets_root: Path,
    brand: TemplateBrand = None, ctx: SlideRenderContext | None = None,
) -> None:
    b = brand or _DEFAULT_BRAND
    ty = ctx.typography if ctx else TypographyScale()
    if not b.has_blueprint_bg:
        _slide_bg(slide, b.white)
    _logo(slide, assets_root, b)
    m = _render_content_header(
        slide, spec.get("title", "Timeline"), b, ctx, ty,
        section_marker=spec.get("section_marker"),
    )
    milestones = spec.get("milestones", [])[:6]
    if not milestones:
        return

    n = len(milestones)
    line_y = m.content_top + int(m.content_height * 0.42)
    line_h = Pt(2)

    # Connector line
    _box(slide, m.margin_l, line_y, m.content_width, line_h,
         fill=_rgb("DDDDDD"))

    segment_w = m.content_width // n
    dot_r = Inches(0.12)
    for i, ms in enumerate(milestones):
        cx = m.margin_l + segment_w * i + segment_w // 2
        # Milestone dot
        dot_shape = slide.shapes.add_shape(
            9, cx - dot_r, line_y - dot_r + line_h // 2,
            dot_r * 2, dot_r * 2,
        )
        _fill_solid(dot_shape, b.primary)
        _no_line(dot_shape)

        label_w = int(segment_w * 0.9)
        label_x = cx - label_w // 2
        # Date above line
        _text_box(
            slide, label_x, line_y - Inches(0.55),
            label_w, Inches(0.4),
            ms.get("date", ""), font_name=b.font_heading,
            font_size=ty.card_headline, bold=True, color=b.primary,
            align=PP_ALIGN.CENTER,
        )
        # Title below line
        _text_box(
            slide, label_x, line_y + Inches(0.2),
            label_w, Inches(0.35),
            ms.get("title", ""), font_name=b.font_heading,
            font_size=12, bold=True, color=b.black,
            align=PP_ALIGN.CENTER,
        )
        # Description
        if ms.get("description"):
            _text_box(
                slide, label_x, line_y + Inches(0.55),
                label_w, Inches(0.6),
                ms["description"], font_name=b.font_body,
                font_size=ty.body_small, color=b.dark_grey,
                align=PP_ALIGN.CENTER, word_wrap=True,
            )
    _source_line(slide, spec.get("source", "Source: [cite]"), b, ctx)


def render_bar_chart(
    slide, spec: dict, assets_root: Path,
    brand: TemplateBrand = None, ctx: SlideRenderContext | None = None,
) -> None:
    b = brand or _DEFAULT_BRAND
    ty = ctx.typography if ctx else TypographyScale()
    if not b.has_blueprint_bg:
        _slide_bg(slide, b.white)
    _logo(slide, assets_root, b)
    m = _render_content_header(
        slide, spec.get("title", "Chart"), b, ctx, ty,
        section_marker=spec.get("section_marker"),
    )
    bars = spec.get("bars", [])[:8]
    if not bars:
        return

    chart_h = int(m.content_height * 0.72)
    chart_top = m.content_top
    max_val = max((bar.get("value", 0) for bar in bars), default=1) or 1

    n = len(bars)
    gap = int(Inches(0.15))
    bar_w = (m.content_width - gap * (n + 1)) // n

    for i, bar in enumerate(bars):
        val = bar.get("value", 0)
        bar_h = max(int(chart_h * val / max_val), Pt(4))
        x = m.margin_l + gap + i * (bar_w + gap)
        y = chart_top + chart_h - bar_h

        color_hex = bar.get("color")
        fill_color = _rgb(color_hex.lstrip("#")) if color_hex else (
            b.primary if i % 2 == 0 else b.dark_grey
        )
        rect = _box(slide, x, y, bar_w, bar_h, fill=fill_color)
        _no_line(rect)

        # Value label above bar
        _text_box(
            slide, x, y - Inches(0.3), bar_w, Inches(0.28),
            str(val), font_name=b.font_heading,
            font_size=11, bold=True, color=b.black,
            align=PP_ALIGN.CENTER,
        )
        # Category label below chart area
        _text_box(
            slide, x, chart_top + chart_h + Inches(0.05),
            bar_w, Inches(0.35),
            bar.get("label", ""), font_name=b.font_body,
            font_size=ty.body_small, color=b.dark_grey,
            align=PP_ALIGN.CENTER, word_wrap=True,
        )

    # Baseline
    _box(slide, m.margin_l, chart_top + chart_h, m.content_width, Pt(1),
         fill=b.dark_grey)
    _source_line(slide, spec.get("source", "Source: [cite]"), b, ctx)


# ---------------------------------------------------------------------------
# Slide dispatcher
# ---------------------------------------------------------------------------

RENDERERS = {
    "title-block":         render_title_block,
    "divider":             render_divider,
    "agenda":              render_agenda,
    "metric-card":         render_metric_cards,
    "challenge-list":      render_challenge_list,
    "tech-tile":           render_tech_tiles,
    "quote-block":         render_quote_block,
    "image-content":       render_image_content,
    "data-table":          render_data_table,
    "recommendation-card": render_recommendation_cards,
    "timeline":            render_timeline,
    "bar-chart":           render_bar_chart,
    "closing":             render_closing,
}


# ---------------------------------------------------------------------------
# Main build function
# ---------------------------------------------------------------------------

def build_deck(
    slides_spec: list[dict],
    assets_root: Path,
    output_path: Path,
    skills_root: Path | None = None,
    template_id: str | None = None,
    master_alias: str | None = None,
) -> Path:
    """Render all slides and save the PPTX.

    When a skeleton.pptx + theme-manifest.yaml exist for the template, the deck
    is built by opening the skeleton and adding slides that INHERIT their
    background from the matched layout (cover/section/content variant).  This
    keeps the real Red Hat master, embedded fonts, and background images intact.

    When no skeleton is found, falls back to the blueprint-replay approach
    (blank Presentation + decorative shapes drawn as explicit shapes).
    """
    # ── Brand ────────────────────────────────────────────────────────────────
    brand = _DEFAULT_BRAND
    tmpl_yaml: Path | None = None
    if skills_root and template_id:
        tmpl_yaml = Path(skills_root) / "templates" / template_id / "template.yaml"
        brand = TemplateBrand.from_template_yaml(tmpl_yaml)
        print(f"  Brand: font={brand.font_heading!r}  primary={brand.primary}")

    # ── Skeleton mode ─────────────────────────────────────────────────────────
    prs: Presentation | None = None
    layout_map: dict = {}
    master_idx = 0

    if skills_root and template_id:
        manifest = _load_theme_manifest(Path(skills_root), template_id)
        if manifest:
            skeleton_rel = manifest.get("skeleton", "skeletons/skeleton.pptx")
            skeleton_path = Path(skills_root) / "templates" / template_id / skeleton_rel
            if skeleton_path.exists():
                prs = Presentation(str(skeleton_path))
                masters_info = manifest.get("masters", {})
                # Resolve master alias (caller override → first in manifest)
                if master_alias and master_alias in masters_info:
                    chosen_alias = master_alias
                else:
                    chosen_alias = next(iter(masters_info), None)
                info = masters_info.get(chosen_alias, {})
                master_idx = info.get("master_index", 0)
                variants = info.get("slide_variants", {})
                layout_map = _build_layout_map(prs, master_idx, variants)
                brand.has_blueprint_bg = True
                print(f"  Skeleton: master='{chosen_alias}'  idx={master_idx}  "
                      f"variants={list(layout_map.keys())}")

    # ── Blueprint fallback (no skeleton) ─────────────────────────────────────
    blueprints: dict = {}
    img_dir: Path | None = None

    if prs is None:
        prs = Presentation()
        prs.slide_width  = SLIDE_W
        prs.slide_height = SLIDE_H
        if skills_root and template_id:
            bp_path = (Path(skills_root) / "templates" / template_id
                       / "layouts" / "blueprints.yaml")
            if bp_path.exists():
                with bp_path.open() as f:
                    blueprints = yaml.safe_load(f) or {}
                img_dir = Path(skills_root) / "assets" / "layout-images"
                brand.has_blueprint_bg = True
                print(f"  Blueprints: {len(blueprints.get('layouts', {}))} layouts  "
                      f"(fallback — no skeleton found)")
            else:
                print("  WARNING: No skeleton or blueprints — slides will lack template decorations")

    # ── Cache slide dimensions ────────────────────────────────────────────────
    prs_width_emu_cache.clear()
    prs_height_emu_cache.clear()
    prs_width_emu_cache.append(int(prs.slide_width))
    prs_height_emu_cache.append(int(prs.slide_height))

    # ── Render slides ─────────────────────────────────────────────────────────
    for i, spec in enumerate(slides_spec):
        element = spec.get("element", "title-block")

        # ── Pick slide base ──────────────────────────────────────────────────
        if layout_map:
            # Skeleton mode: inherit background from the variant-matched layout
            variant = _ELEMENT_VARIANT.get(element, "content")
            layout  = layout_map.get(variant, next(iter(layout_map.values())))
            slide   = prs.slides.add_slide(layout)
            # Read placeholder zones BEFORE clearing (they define content area)
            ph_zones = _read_placeholder_zones(
                slide, prs_width_emu_cache[0], prs_height_emu_cache[0],
            )
            _clear_placeholders(slide)
            ctx = _ctx_from_placeholders(
                ph_zones, prs_width_emu_cache[0], prs_height_emu_cache[0],
                template_id, tmpl_yaml,
            )
            ctx.dark_bg = variant in ("cover", "section")
        else:
            # Blueprint fallback
            layout_name = spec.get("layout_name") or _element_to_layout(
                element,
                skills_root=Path(skills_root) if skills_root else None,
                template_id=template_id,
            )
            if layout_name and template_id and str(template_id).startswith("upload-"):
                layout_name = resolve_slide_layout_name(prs, layout_name, template_id)
            bp_layout = (blueprints.get("layouts") or {}).get(layout_name or "")
            slide_master_idx = master_idx
            if bp_layout:
                slide_master_idx = bp_layout.get("master", master_idx)
            if slide_master_idx >= len(prs.slide_masters):
                slide_master_idx = 0
            slide = _add_empty_slide(prs, slide_master_idx, layout_name_hint="blank")
            _slide_bg(slide, WHITE)
            if blueprints and img_dir and layout_name:
                _replay_blueprints(slide, layout_name, blueprints, img_dir)
            ctx = None
            if blueprints and layout_name:
                ctx = SlideRenderContext.from_blueprints(
                    blueprints, layout_name,
                    prs_width_emu_cache[0], prs_height_emu_cache[0],
                    template_id, template_yaml_path=tmpl_yaml,
                )
                variant = _ELEMENT_VARIANT.get(element, "content")
                ctx.dark_bg = _detect_dark_bg(
                    bp_layout, img_dir,
                    fallback=variant in ("section",),
                )

        # ── Render content ───────────────────────────────────────────────────
        renderer = RENDERERS.get(element)
        if renderer:
            try:
                renderer(slide, spec, assets_root, brand, ctx)
            except Exception as exc:
                _text_box(
                    slide, MARGIN_L, Inches(3), CONTENT_W, Inches(1),
                    f"[Slide {i+1}: {element}] Render error: {exc}",
                    font_name=brand.font_body, font_size=14, color=_rgb("CC0000"),
                )
        else:
            _text_box(
                slide, MARGIN_L, Inches(3), CONTENT_W, Inches(1),
                f"[Unknown element: {element}]",
                font_name=brand.font_body, font_size=14, color=brand.mid_grey,
            )

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"  Saved: {output_path}  ({output_path.stat().st_size // 1024} KB)")
    return output_path


def _element_to_layout(
    element: str,
    skills_root: Path | None = None,
    template_id: str | None = None,
) -> str | None:
    """Map element type to blueprint layout key.

    Loads from element YAML if skills_root is provided, falling back to hardcoded map.
    """
    if skills_root and template_id:
        elem_yaml = Path(skills_root) / "templates" / template_id / "elements" / f"{element}.yaml"
        if elem_yaml.exists():
            with elem_yaml.open() as f:
                data = yaml.safe_load(f) or {}
            layout = data.get("blueprint_layout")
            if layout:
                return layout

    # Hardcoded fallback
    mapping = {
        # Master 1 layouts (teal wave / Sales Enablement brand)
        "title-block":         "M1:Front cover-02",
        "closing":             "M1:Blank page_1",  # overridden per template in element YAML
        "divider":             "M1:Front cover-02_1",
        "agenda":              "M1:Text page",
        "metric-card":         "M1:Text page",
        "challenge-list":      "M1:Text page",
        "tech-tile":           "M1:Text page",
        "image-content":       "M1:Text page_2",
        "quote-block":         "M1:Text page_2",
        "recommendation-card": "M1:Text page",
        "data-table":          "M1:Text page",
        "timeline":            "M1:Text page",
    }
    return mapping.get(element)
